#!/usr/bin/env python3
"""
Dédalo — the curator's presentation (30-minute version, 28 slides).

Self-contained and self-authored: everything is drawn with native PowerPoint
shapes, so every slide stays editable and nothing is a raster image.

ONE source for both artifacts: this file builds `dedalo_for_curators.pptx`
(slides + embedded speaker notes) and `speaker_script.md` (the same notes,
printable). Edit here and re-run; never edit the .pptx by hand and expect the
script to follow.

    python -m venv .venv && .venv/bin/pip install python-pptx
    .venv/bin/python presentation/curators_30min/build_deck.py

This directory is INDEPENDENT of presentation/build_deck.py, which is a
different deck maintained separately.
"""
import math
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
W, H = 13.333, 7.5
M = 0.9                      # page margin
COL = W - 2 * M              # content width

# ── palette ────────────────────────────────────────────────────────────────
INK      = C(0x1A, 0x1C, 0x22)
PAPER    = C(0xFA, 0xF8, 0xF4)
DEEP     = C(0x14, 0x1A, 0x24)
ORANGE   = C(0xC2, 0x62, 0x2A)
ORANGE_L = C(0xE8, 0x9A, 0x3C)
TEAL     = C(0x0E, 0x74, 0x90)
GREEN    = C(0x15, 0x80, 0x3D)
VIOLET   = C(0x6D, 0x28, 0xD9)
ROSE     = C(0xBE, 0x12, 0x3C)
MUTED    = C(0x66, 0x6E, 0x7A)
MUTED_D  = C(0x93, 0x9E, 0xAC)
HAIR     = C(0xE2, 0xDC, 0xD1)
WHITE    = C(0xFF, 0xFF, 0xFF)
SLATE    = C(0x47, 0x55, 0x69)
TINT = {ORANGE: C(0xFB, 0xF1, 0xE7), TEAL: C(0xEC, 0xF6, 0xF9),
        GREEN:  C(0xED, 0xF8, 0xF1), VIOLET: C(0xF4, 0xEF, 0xFD),
        ROSE:   C(0xFD, 0xEF, 0xF2), SLATE: C(0xF2, 0xF0, 0xEC),
        MUTED:  C(0xF2, 0xF0, 0xEC)}

FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]
SCRIPT = []


# ── primitives ─────────────────────────────────────────────────────────────
def slide(bg=PAPER, title_for_script=""):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    s._script_title = title_for_script
    return s


def notes(s, text_):
    s.notes_slide.notes_text_frame.text = text_.strip()
    SCRIPT.append((getattr(s, "_script_title", ""), text_.strip()))


def _spc(run, hundredths_pt):
    run.font._element.set("spc", str(hundredths_pt))     # rPr/@spc, 1/100 pt


def est_lines(content, w_in, size):
    """Rough wrapped-line count. Helvetica Neue averages a little over half an
    em per character; renderers differ, so this errs on the generous side."""
    cpl = max(8, int(w_in * 72 / (size * 0.56)))
    items = content if isinstance(content, list) else [content]
    n = 0
    for it in items:
        t = it[0] if isinstance(it, tuple) else it
        n += max(1, -(-len(t) // cpl))
    return n


def line_h(size, spacing=1.3):
    return size * spacing / 72.0


def text(s, x, y, w, h, lines, size=14, color=MUTED, bold=False, align=PP_ALIGN.LEFT,
         spacing=1.25, after=0, tracking=None, italic=False):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        over = {}
        if isinstance(ln, tuple):
            ln, over = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.line_spacing = over.get("spacing", spacing)
        p.space_after = Pt(over.get("after", after))
        r = p.add_run()
        r.text = ln
        r.font.name = FONT
        r.font.size = Pt(over.get("size", size))
        r.font.bold = over.get("bold", bold)
        r.font.italic = over.get("italic", italic)
        r.font.color.rgb = over.get("color", color)
        tr = over.get("tracking", tracking)
        if tr:
            _spc(r, tr)
    return tb


def kicker(s, x, y, label, color=ORANGE, size=11.5):
    return text(s, x, y, COL, 0.3, label.upper(), size=size, color=color, bold=True,
                tracking=160)


def title(s, txt, x=M, y=0.95, w=None, size=34, color=INK, spacing=1.08):
    return text(s, x, y, w or COL, 1.6, txt, size=size, color=color, bold=True,
                spacing=spacing)


def rule(s, x, y, w, color=HAIR, thickness=0.014):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y), In(w), In(thickness))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def box(s, x, y, w, h, fill=WHITE, line=HAIR, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.045, line_w=1.0):
    sh = s.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    sh.text_frame.text = ""
    return sh


def accent_bar(s, x, y, h, color, w=0.055):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.5
    except Exception:
        pass
    return sh


def disc(s, cx, cy, d, fill, label="", color=WHITE, size=18):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, In(cx - d / 2), In(cy - d / 2), In(d), In(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    if label:
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = color
    return sh


def shape_icon(s, kind, cx, cy, size, color):
    presets = {"db": MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
               "doc": MSO_SHAPE.FLOWCHART_DOCUMENT,
               "docs": MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
               "screen": MSO_SHAPE.FLOWCHART_DISPLAY,
               "cloud": MSO_SHAPE.CLOUD}
    sh = s.shapes.add_shape(presets[kind], In(cx - size / 2), In(cy - size / 2),
                            In(size), In(size * 0.82))
    sh.fill.background()
    sh.line.color.rgb = color; sh.line.width = Pt(1.9)
    sh.shadow.inherit = False
    return sh


def arrow_h(s, x, y, length, color=C(0x9A, 0xA4, 0xB2), thick=0.028, head=0.16):
    stem = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y - thick / 2),
                              In(length - head * 0.75), In(thick))
    stem.fill.solid(); stem.fill.fore_color.rgb = color
    stem.line.fill.background(); stem.shadow.inherit = False
    tip = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, In(x + length - head),
                             In(y - head / 2), In(head), In(head))
    tip.fill.solid(); tip.fill.fore_color.rgb = color
    tip.line.fill.background(); tip.shadow.inherit = False
    tip.rotation = 90
    return stem


def card(s, x, y, w, h, heading, lines, color=ORANGE, tint=False, num=None,
         hsize=18, bsize=13):
    """A card whose body starts BELOW its heading, however many lines that took."""
    box(s, x, y, w, h, fill=(TINT[color] if tint else WHITE))
    accent_bar(s, x, y + 0.28, h - 0.56, color)
    tx, ty = x + 0.34, y + 0.3
    if num is not None:
        disc(s, tx + 0.24, ty + 0.22, 0.48, color, str(num), size=16)
        ty += 0.72
    hl = est_lines(heading, w - 0.62, hsize)
    text(s, tx, ty, w - 0.62, hl * line_h(hsize, 1.15) + 0.1, heading, size=hsize,
         color=INK, bold=True, spacing=1.12)
    if lines:
        text(s, tx, ty + hl * line_h(hsize, 1.15) + 0.18, w - 0.62, h - 1.0, lines,
             size=bsize, color=MUTED, spacing=1.32, after=3)


def footer(s, txt="Dédalo · Cultural Heritage management", page=None):
    grey = C(0xA6, 0x9E, 0x92)
    text(s, M, H - 0.6, COL * 0.7, 0.3, txt, size=10, color=grey, tracking=60)
    if page is not None:
        text(s, W - M - 1.2, H - 0.6, 1.2, 0.3, str(page), size=10, color=grey,
             align=PP_ALIGN.RIGHT)


def divider(part, heading, sub, color=ORANGE, script_title=""):
    s = slide(DEEP, script_title or heading.replace("\n", " "))
    nl = heading.count("\n") + 1
    accent_bar(s, M, 2.5, 1.75 + nl * 0.8, color, w=0.07)
    text(s, M + 0.45, 2.5, COL, 0.4, part.upper(), size=12, color=ORANGE_L, bold=True,
         tracking=200)
    text(s, M + 0.45, 3.0, COL * 0.8, nl * 0.85, heading, size=42, color=WHITE,
         bold=True, spacing=1.14)
    text(s, M + 0.45, 3.0 + nl * 0.78 + 0.24, COL * 0.6, 0.7, sub, size=17,
         color=MUTED_D, spacing=1.35)
    return s


# ═══ 1 · TITLE ═════════════════════════════════════════════════════════════
s = slide(DEEP, "Title")
box(s, 0, 0, W, 0.34, fill=ORANGE, line=None, shape=MSO_SHAPE.RECTANGLE)
text(s, M, 2.15, COL, 0.4, "FREE AND OPEN SOURCE · SINCE 2003", size=12,
     color=ORANGE_L, bold=True, tracking=200)
text(s, M, 2.72, COL, 1.5, "Dédalo", size=88, color=WHITE, bold=True, spacing=1.0)
text(s, M, 4.12, COL * 0.72, 1.2, "Cultural heritage, catalogued to last.",
     size=30, color=C(0xD8, 0xDF, 0xE8), spacing=1.2)
rule(s, M, 5.35, 2.2, ORANGE, 0.03)
text(s, M, 5.7, COL * 0.7, 0.8,
     ["A thirty-minute introduction for curators, archivists and collection teams",
      "No technical background needed."], size=14, color=MUTED_D, spacing=1.5)
notes(s, """
Good morning, and thank you for the time.

For the next half hour I want to show you what Dédalo is and, more importantly, what it
is FOR. There is almost no technology in this talk. Nearly everything I am going to say is
about how a catalogue behaves over twenty years, because that is the real problem in our
field — not storing records, but keeping them coherent, correct and usable long after the
people who made them have moved on.

Three things up front. Dédalo is free and open source: there is no licence to buy and no
seat to pay for. It has been in production in real institutions for over twenty years —
archives, museums, oral memory projects, archaeology, numismatics. And its engine was
rebuilt from scratch recently, so what I am showing you is a mature idea on new
foundations.

Please interrupt with questions whenever you like. It works better as a conversation.
""")

# ═══ 2 · HOOK ══════════════════════════════════════════════════════════════
s = slide(DEEP, "Hook — our archives know more than our searches can find")
accent_bar(s, M, 1.6, 3.9, ORANGE, w=0.07)
text(s, M + 0.45, 1.75, COL - 0.6, 3.4,
     ["Our archives know more", "than our searches can find."],
     size=52, color=WHITE, bold=True, spacing=1.14)
text(s, M + 0.5, 5.4, COL * 0.66, 1.0,
     "Hold that sentence. We come back to it at the end.", size=16, color=MUTED_D)
notes(s, """
Let me start with a true story, because it is the shape of the whole problem.

An oral history archive holds ten thousand interviews about life in a rural valley. A
researcher wants every testimony that touches on people displaced when the reservoir was
built. She types "reservoir" into the search box. She gets the handful of interviews where
somebody happened to use that word.

She misses the ones that say "when the water came", "they flooded our houses", "el
pantano", "when they made us leave". Dozens of them. The knowledge is IN the archive. The
search cannot reach it, because a classic search matches strings of letters, not meaning.

That gap — between what an archive knows and what it can be asked — is what a good heritage
system should be closing. Keep that sentence in mind: we come back to it at the very end,
and I will show you what Dédalo does about it.
""")

# ═══ 3 · WHAT IT IS ════════════════════════════════════════════════════════
s = slide(PAPER, "What Dédalo is")
kicker(s, M, 0.72, "In one sentence")
title(s, "A system for managing cultural heritage —\nbuilt to produce good data.", y=1.12)
rule(s, M, 2.65, COL)
items = [("Free and open source", "No licence, no seat fee, no vendor\nlock. The code and your data are yours.", ORANGE),
         ("For every discipline", "Archives, museums, oral memory,\narchaeology, ethnology, numismatics.", TEAL),
         ("Twenty years in production", "Designed with institutions, in\ninstitutions — not in a start-up.", GREEN),
         ("Rebuilt, not patched", "A brand-new engine under a mature\nmodel, ready for the next twenty years.", VIOLET)]
cw = (COL - 3 * 0.28) / 4
for i, (h, b, col) in enumerate(items):
    card(s, M + i * (cw + 0.28), 3.0, cw, 2.8, h, b.split("\n"), col, tint=True,
         hsize=16, bsize=12.5)
footer(s, page=3)
notes(s, """
So, in one sentence: Dédalo is a system for managing cultural heritage that is built to
produce good data.

That last phrase is deliberate. Plenty of software will store what you type. Very little of
it is designed around the question "will this still be coherent, correct and reusable in
twenty years?" That is the question Dédalo is organised around.

Four things worth knowing. It is free and open source — no licence, no per-seat fee, and
crucially no lock-in: your data is in an open database on a machine you control. It is not
built for one discipline; the same system runs an oral memory archive, a coin cabinet and
an archaeological inventory, and I will explain in a moment how that is possible without it
becoming a generic mush.

It has been in production for two decades, developed with institutions rather than pitched
at them. And its engine was recently rebuilt from the ground up — so the model is mature,
and the foundation is new.
""")

# ═══ 4 · NOT A CMS ═════════════════════════════════════════════════════════
s = slide(PAPER, "Not a content manager")
kicker(s, M, 0.72, "A common misunderstanding")
title(s, "It is not a website manager\nthat happens to hold records.", y=1.12)
lw = (COL - 0.5) / 2
box(s, M, 3.05, lw, 2.7, fill=TINT[SLATE], line=HAIR)
text(s, M + 0.42, 3.35, lw - 0.84, 0.4, "A CONTENT MANAGER", size=11.5, color=MUTED,
     bold=True, tracking=160)
text(s, M + 0.42, 3.85, lw - 0.9, 1.6,
     ["Helps you publish pages.", "The page is the point; the data exists to fill it."],
     size=17, color=MUTED, spacing=1.35, after=8)
box(s, M + lw + 0.5, 3.05, lw, 2.7, fill=TINT[ORANGE], line=C(0xEE, 0xD6, 0xBE))
accent_bar(s, M + lw + 0.5, 3.33, 2.14, ORANGE)
text(s, M + lw + 0.92, 3.35, lw - 0.84, 0.4, "DÉDALO", size=11.5, color=ORANGE,
     bold=True, tracking=160)
text(s, M + lw + 0.92, 3.85, lw - 1.0, 1.6,
     ["Helps you produce data.", "The data is the point; the page is one of its outputs."],
     size=17, color=INK, spacing=1.35, after=8)
footer(s, page=4)
notes(s, """
The most common misunderstanding I meet is that Dédalo is a kind of website manager for
museums. It is worth clearing that up early, because it changes what you should expect
from it.

A content manager — WordPress and its cousins — helps you publish pages. The page is the
point. The data exists to fill the page, and it usually lives inside the page: a
description typed here, the same place name typed again there, a date written three
different ways because three different people wrote it.

Dédalo works the other way round. The data is the point. The website is simply one of the
things your data can be turned into — along with an export for a national aggregator, a
research dataset, a printed catalogue, or a partner's portal.

That inversion is not a slogan. It has consequences on every screen your team will touch,
and the rest of this talk is really about those consequences.
""")

# ═══ 5 · PART ONE ══════════════════════════════════════════════════════════
s = divider("Part one", "How a catalogue is built",
            "Three words, and one idea that changes everything.")
notes(s, """
Part one. How a catalogue is actually built in Dédalo.

I need to give you three words — words you already have, just with precise meanings here —
and then one single idea. If you take only one thing away from this whole talk it will be
that idea, and it is coming up in about two minutes.
""")

# ═══ 6 · THREE WORDS ═══════════════════════════════════════════════════════
s = slide(PAPER, "Three words")
kicker(s, M, 0.72, "The vocabulary")
title(s, "Three words cover most of what you do.", y=1.12)
cw = (COL - 2 * 0.34) / 3
trio = [("Record", "One thing you catalogue.", "An object. An interview.\nA photograph. A person. A place.", "doc", ORANGE),
        ("Section", "Records of the same kind.", "“Coins”. “Oral testimonies”.\n“People”. A series of like units.", "docs", TEAL),
        ("Component", "One field on the record.", "A title. A date. An extent.\nA photograph. A link.", "screen", VIOLET)]
for i, (h, sub, b, icon, col) in enumerate(trio):
    x = M + i * (cw + 0.34)
    box(s, x, 2.6, cw, 3.5)
    accent_bar(s, x, 2.9, 2.9, col)
    shape_icon(s, icon, x + 0.75, 3.25, 0.66, col)
    text(s, x + 0.4, 3.9, cw - 0.8, 0.55, h, size=26, color=INK, bold=True)
    text(s, x + 0.4, 4.48, cw - 0.8, 0.4, sub, size=15, color=col, bold=True)
    text(s, x + 0.4, 5.0, cw - 0.8, 1.1, b.split("\n"), size=13, color=MUTED,
         spacing=1.35, after=3)
footer(s, page=6)
notes(s, """
Three words. You know all three concepts already; I just want us to use the same names.

A RECORD is one thing you catalogue. An object, an interview, a photograph, a person, a
place, a bibliographic reference.

A SECTION holds all the records of the same kind. "Coins" is a section. "Oral testimonies"
is a section. "People" is a section. If you come from databases, a section is a table. If
you come from archives, think of a series of homogeneous units of description.

And a COMPONENT is one field on the record — a title, a date, an extent, a photograph, a
link to another record.

One detail that matters more than it sounds: components have KINDS, and the kind makes the
field intelligent. A date field understands imperfect and approximate dates. A geographic
field understands coordinates. A media field knows a file has a master and derivatives. A
link field knows its content is a pointer to another record, not a piece of text.
""")

# ═══ 7 · THE BIG IDEA ══════════════════════════════════════════════════════
s = slide(DEEP, "THE BIG IDEA — the shape of your catalogue is data")
text(s, M, 0.85, COL, 0.4, "THE IDEA WORTH REMEMBERING", size=12, color=ORANGE_L,
     bold=True, tracking=200)
text(s, M, 1.32, COL * 0.95, 1.5, "The shape of your catalogue is data.",
     size=44, color=WHITE, bold=True, spacing=1.1)
lw = (COL - 0.55) / 2
box(s, M, 3.3, lw, 2.95, fill=C(0x21, 0x2A, 0x36), line=C(0x33, 0x3E, 0x4C))
text(s, M + 0.45, 3.62, lw - 0.9, 0.4, "IN MOST SYSTEMS", size=11.5, color=MUTED_D,
     bold=True, tracking=160)
text(s, M + 0.45, 4.1, lw - 0.9, 1.4,
     ["The shape lives in the program.",
      "Which record types exist, which fields, in what order, in which languages — written in code."],
     size=15, color=C(0xB9, 0xC3, 0xD0), spacing=1.35, after=8)
text(s, M + 0.45, 5.72, lw - 0.9, 0.4, "Changing it = a developer, a release, a budget.",
     size=13.5, color=ROSE, bold=True)
box(s, M + lw + 0.55, 3.3, lw, 2.95, fill=C(0x2A, 0x22, 0x1A), line=ORANGE, line_w=1.4)
text(s, M + lw + 1.0, 3.62, lw - 0.9, 0.4, "IN DÉDALO", size=11.5, color=ORANGE_L,
     bold=True, tracking=160)
text(s, M + lw + 1.0, 4.1, lw - 0.9, 1.4,
     ["The shape is data you edit.",
      "Record types, fields, order, vocabularies, languages — described inside Dédalo."],
     size=15, color=C(0xEE, 0xE6, 0xDC), spacing=1.35, after=8)
text(s, M + lw + 1.0, 5.72, lw - 0.9, 0.4, "Changing it = an act of documentation.",
     size=13.5, color=ORANGE_L, bold=True)
notes(s, """
Here is the idea. If you remember nothing else, remember this.

In most systems, the shape of your catalogue lives inside the program. Which record types
exist, which fields they have, in what order, with which vocabularies, in which languages —
all of that is written in code by a developer. Which means that changing it requires a
developer, a release, a budget and a delay. Every curator in this room has waited months
for a field.

In Dédalo, the shape of your catalogue is itself DATA. It is described inside Dédalo and
edited inside Dédalo, like any other content. Adding a field, creating a whole new type of
record, reordering a form, making a field repeatable or translatable, saying "this field
must draw its values from that thesaurus" — those are acts of documentation, not acts of
programming.

The system calls that description an ontology. Do not let the word frighten anyone: it
simply means the catalogue knows what it is made of.
""")

# ═══ 8 · WHAT THAT CHANGES ═════════════════════════════════════════════════
s = slide(PAPER, "What that changes")
kicker(s, M, 0.72, "In practice")
title(s, "What that changes, on a Tuesday afternoon.", y=1.12)
rows = [("Add a field to a record type", "A ticket, a release, a wait", "You do it, and catalogue with it today"),
        ("Create a new type of record", "A project, and a quotation", "A decision, taken in a meeting"),
        ("A standard or vocabulary changes", "A dependency on your supplier", "Your own work, on your own schedule"),
        ("Publish to a new aggregator", "A development contract", "A configuration, then publish again")]
y = 2.6
box(s, M, y - 0.22, COL, 0.62, fill=TINT[SLATE], line=None)
text(s, M + 0.36, y - 0.04, 4.5, 0.35, "WHEN YOU NEED TO…", size=11, color=MUTED,
     bold=True, tracking=150)
text(s, M + 5.2, y - 0.04, 3.2, 0.35, "USUALLY", size=11, color=MUTED, bold=True, tracking=150)
text(s, M + 8.4, y - 0.04, 3.2, 0.35, "IN DÉDALO", size=11, color=ORANGE, bold=True, tracking=150)
y += 0.62
for i, (a, b, c) in enumerate(rows):
    text(s, M + 0.36, y + 0.14, 4.6, 0.4, a, size=15, color=INK, bold=True)
    text(s, M + 5.2, y + 0.16, 3.1, 0.4, b, size=13.5, color=MUTED)
    text(s, M + 8.4, y + 0.16, 3.4, 0.4, c, size=13.5, color=ORANGE, bold=True)
    if i < len(rows) - 1:
        rule(s, M + 0.36, y + 0.66, COL - 0.72)
    y += 0.78
box(s, M, 6.22, COL, 0.56, fill=TINT[TEAL], line=None)
text(s, M + 0.36, 6.39, COL - 0.7, 0.35,
     "The honest part: this power sits with whoever designs your model.",
     size=13.5, color=TEAL, bold=True)
footer(s, page=8)
notes(s, """
What does that actually change? Let me be concrete, because this is where the idea stops
being abstract.

You need a new field on a record type. Usually: a ticket to your supplier, a release, a
wait of weeks or months. Here: you add it, and your team catalogues with it the same
afternoon.

You need a whole new kind of record — say you start accepting audiovisual donations.
Usually a project with a quotation. Here, a decision taken in a meeting.

A standard changes, or your thesaurus authority publishes a new version. Usually you depend
on your supplier's roadmap. Here it is your own work, on your own schedule. And a new
aggregator wants your data in their profile: configuration, then publish again.

Now the honest part, and I want to say it out loud. This power sits with whoever designs
your model. Dédalo will not stop you designing a bad catalogue — it stops the software from
imposing one on you. So treat the model design with the same care as a cataloguing manual.
""")


# ═══ 9 · CONTROLLED VOCABULARY ═════════════════════════════════════════════
s = slide(PAPER, "Words that mean something")
kicker(s, M, 0.72, "Controlled vocabularies")
title(s, "A term is a record too.", y=1.12, size=36)
text(s, M, 2.0, COL * 0.5, 0.9,
     "Places, people, subjects, materials — held as a tree, and every term in it is a full record with its own description, variants and history.",
     size=15, color=MUTED, spacing=1.4)
tx, ty = M + 0.1, 3.45
box(s, tx, ty, 2.5, 0.62, fill=TINT[TEAL], line=C(0xCF, 0xE4, 0xEC))
text(s, tx + 0.3, ty + 0.19, 2.1, 0.3, "Toponymy", size=15, color=TEAL, bold=True)
rule(s, tx + 0.42, ty + 0.62, 0.014, C(0xC6, 0xCE, 0xD8), 2.05)
for i, k in enumerate(["Alt Urgell", "Vall de Cabó", "La Seu d'Urgell"]):
    yy = ty + 0.95 + i * 0.72
    rule(s, tx + 0.42, yy + 0.3, 0.45, C(0xC6, 0xCE, 0xD8), 0.012)
    box(s, tx + 0.9, yy, 2.6, 0.6)
    text(s, tx + 1.15, yy + 0.18, 2.3, 0.3, k, size=14, color=INK, bold=(i == 1))
box(s, M + 6.3, 3.05, COL - 6.3, 3.1, fill=WHITE)
accent_bar(s, M + 6.3, 3.33, 2.54, TEAL)
text(s, M + 6.7, 3.32, 4.4, 0.4, "VALL DE CABÓ — ONE RECORD", size=11.5, color=TEAL,
     bold=True, tracking=140)
yy = 3.9
for k, v in [("Preferred term", "Vall de Cabó"), ("Variants", "Vall de Cabo · Valle de Cabó"),
             ("Coordinates", "42.19 N, 1.36 E"), ("Described in", "ca · es · en"),
             ("Referred to by", "914 records")]:
    text(s, M + 6.7, yy, 1.9, 0.3, k, size=12.5, color=MUTED)
    text(s, M + 8.75, yy, 2.9, 0.3, v, size=12.5, color=INK, bold=True)
    yy += 0.44
footer(s, page=9)
notes(s, """
Controlled vocabularies. Every institution has them: place names, personal names, subject
headings, materials, techniques, typologies.

In Dédalo a thesaurus is a tree of terms — and here is the part that matters — the terms in
it are ordinary records. That sounds like a technicality. It is not.

It means a place in your toponymic thesaurus is not the string "Vall de Cabó" typed into
nine hundred records. It is ONE record. It has its own preferred form and its variants. It
has coordinates. It is described in each of your working languages. It can carry its own
history, its own bibliography, its own note about why the spelling changed in 1934.

And nine hundred records point at it.

That is the difference between a vocabulary you maintain and a vocabulary you merely hope
everybody types the same way.
""")

# ═══ 10 · LINKS NOT COPIES ═════════════════════════════════════════════════
s = slide(PAPER, "Links, not copies")
kicker(s, M, 0.72, "The single most important behaviour")
title(s, "Links, not copies.", y=1.12, size=38)
cx, cy = M + 3.1, 4.2
for i in range(10):
    a = math.radians(i * 36)
    disc(s, cx + 2.25 * math.cos(a), cy + 1.5 * math.sin(a), 0.42, C(0xD7, 0xDD, 0xE5))
disc(s, cx, cy, 1.55, TEAL)
text(s, cx - 0.74, cy - 0.34, 1.48, 0.75, ["ONE", "record"], size=15, color=WHITE,
     bold=True, align=PP_ALIGN.CENTER, spacing=1.2)
yy = 2.45
for h, b in [("Correct once.", "Fix the term record and every record that points at it is correct — immediately, everywhere."),
             ("Ask backwards.", "“Which testimonies mention this person?” — a question you never had to plan for."),
             ("Never diverge.", "Two records cannot hold two spellings of one place, because they hold no spelling at all.")]:
    text(s, M + 6.9, yy, 4.9, 0.35, h, size=20, color=INK, bold=True)
    text(s, M + 6.9, yy + 0.4, 4.9, 0.8, b, size=13.5, color=MUTED, spacing=1.35)
    yy += 1.28
box(s, M + 6.9, 6.28, 4.9, 0.55, fill=TINT[ORANGE], line=None)
text(s, M + 7.15, 6.44, 4.5, 0.3, "A catalogue, not a spreadsheet.", size=13.5,
     color=ORANGE, bold=True)
footer(s, page=10)
notes(s, """
This is, for my money, the single most important behaviour in the whole system, and it
follows directly from the previous slide.

When a record refers to another record, Dédalo stores a POINTER — never a copy of its name.
Three consequences, all of them things documentalists want.

First: correct once. You fix a name, a date, a spelling in the term record, and every
record that refers to it is correct. Immediately. Everywhere. No find-and-replace, no
"which of the four spellings is right", no cleanup project.

Second: you can ask backwards. Because the link is real, you can ask "which testimonies
mention this person?" or "what came out of this excavation?" — questions nobody had to
anticipate when the form was designed.

Third: silent divergence becomes impossible. Two records cannot end up with two spellings
of the same place, because neither of them holds a spelling at all.

That is what separates a catalogue from a spreadsheet, and it is why a Dédalo collection is
still usable after twenty years.
""")

# ═══ 11 · LANGUAGES ════════════════════════════════════════════════════════
s = slide(PAPER, "Many languages, two kinds")
kicker(s, M, 0.72, "Multilingual by design")
title(s, "Two things get translated.\nConfusing them causes trouble.", y=1.12)
box(s, M, 3.0, COL, 1.55, fill=TINT[VIOLET], line=C(0xE0, 0xD5, 0xF8))
accent_bar(s, M, 3.26, 1.03, VIOLET)
text(s, M + 0.45, 3.22, 3.0, 0.4, "THE DATA", size=11.5, color=VIOLET, bold=True, tracking=160)
text(s, M + 0.45, 3.66, 6.3, 0.8,
     "A title, a description, an abstract — in as many languages as you work in, each stored on the same record.",
     size=15, color=INK, spacing=1.32)
text(s, M + 7.4, 3.66, 4.2, 0.8,
     "Some values are the same in every language: an inventory number, a coordinate, a year.",
     size=13, color=MUTED, spacing=1.32)
box(s, M, 4.85, COL, 1.55, fill=TINT[SLATE], line=HAIR)
accent_bar(s, M, 5.11, 1.03, MUTED)
text(s, M + 0.45, 5.07, 3.0, 0.4, "THE INTERFACE", size=11.5, color=MUTED, bold=True, tracking=160)
text(s, M + 0.45, 5.51, 6.3, 0.8,
     "The buttons, menus and field labels your team sees while cataloguing.",
     size=15, color=INK, spacing=1.32)
text(s, M + 7.4, 5.51, 4.2, 0.8,
     "A cataloguer can work in Català while entering a description in English.",
     size=13, color=MUTED, spacing=1.32)
footer(s, page=11)
notes(s, """
Multilingualism. Two different things get translated in a heritage system, and confusing
them causes real trouble — I have seen projects lose a year to it.

The first is THE DATA. A title, a description, an abstract can exist in as many languages
as your institution works in, stored separately on the same record. And some values are
deliberately NOT translatable: an inventory number, a coordinate, a year are the same in
every language, and the system knows that.

The second is THE INTERFACE — the buttons, the menus, the field labels your cataloguer
sees.

They are independent. A cataloguer can work with the interface in Català while entering a
description in English, and a researcher can read the catalogue in one language while the
data is served in another.

For those of us in bilingual or trilingual regions, and for anyone who wants their
collection readable beyond their own language, this is not a nice-to-have. It is the
difference between one audience and five.
""")

# ═══ 12 · VERSIONS ═════════════════════════════════════════════════════════
s = slide(PAPER, "Nothing is silently lost")
kicker(s, M, 0.72, "Every change is kept")
title(s, "Nothing is silently lost.", y=1.12, size=38)
rule(s, M + 0.3, 3.6, COL - 1.6, C(0xD7, 0xDD, 0xE5), 0.02)
stops = [("2019", "created", "M. Ferrer"), ("2021", "date corrected", "A. Roca"),
         ("2023", "attribution revised", "M. Ferrer"), ("2024", "linked to the dig", "J. Pla"),
         ("today", "you are here", "")]
for i, (yr, what, who) in enumerate(stops):
    x = M + 0.3 + i * ((COL - 1.6) / (len(stops) - 1))
    last = i == len(stops) - 1
    disc(s, x, 3.61, 0.42 if last else 0.3, ORANGE if last else C(0x9A, 0xA4, 0xB2))
    text(s, x - 1.0, 2.95, 2.0, 0.3, yr, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, x - 1.15, 4.0, 2.3, 0.6, [what] + ([who] if who else []), size=12.5,
         color=MUTED, align=PP_ALIGN.CENTER, spacing=1.3)
box(s, M, 5.3, COL, 1.15, fill=TINT[ORANGE], line=None)
text(s, M + 0.45, 5.55, COL - 1.0, 0.8,
     ["Not a comfort feature — an evidential one.",
      "It makes an attribution defensible, a correction traceable, a disputed record reconstructible."],
     size=15.5, color=INK, bold=True, spacing=1.4, after=6)
footer(s, page=12)
notes(s, """
Every change to a record is kept. What it was before, who changed it, and when. You can
read the history of a record and restore an earlier state.

I want to be careful about how I sell this, because "version history" sounds like a
convenience feature — an undo button.

For a scientific catalogue it is an evidential feature. It is what makes an attribution
DEFENSIBLE: you can show when it changed and on whose authority. It is what makes a
correction TRACEABLE. It is what lets a disputed record be reconstructed years later. And
it is what lets you hand a collection to a colleague without losing the reasoning that
produced it.

It also removes a specific fear that slows teams down: the fear of a bad bulk edit. When a
mistaken import or a wrong batch change is always recoverable, people stop being timid with
their own catalogue.
""")

# ═══ 13 · PERMISSIONS ══════════════════════════════════════════════════════
s = slide(PAPER, "Who may see what")
kicker(s, M, 0.72, "Permissions")
title(s, "Not “editor or reader”.\nField by field.", y=1.12)
cw = (COL - 2 * 0.34) / 3
for i, (h, b, col) in enumerate([("None", "cannot see it exists", ROSE),
                                 ("Read", "can consult, not change", TEAL),
                                 ("Read & write", "can consult and edit", GREEN)]):
    x = M + i * (cw + 0.34)
    box(s, x, 2.8, cw, 1.1, fill=TINT[col], line=None)
    text(s, x + 0.38, 2.99, cw - 0.7, 0.3, h, size=19, color=col, bold=True)
    text(s, x + 0.38, 3.35, cw - 0.7, 0.3, b, size=13, color=MUTED)
text(s, M, 4.25, COL, 0.4, "SET PER SECTION AND PER FIELD — SO THIS IS ORDINARY:",
     size=11.5, color=MUTED, bold=True, tracking=150)
yy = 4.72
for who, can, cannot, col in [("A volunteer", "catalogues the physical description…", "…and never sees the donor's name.", GREEN),
                              ("An external researcher", "reads a collection in full…", "…and cannot alter a character of it.", TEAL),
                              ("A project team", "sees its own material…", "…invisible to the rest of the house.", VIOLET)]:
    box(s, M, yy, COL, 0.58, fill=WHITE)
    accent_bar(s, M, yy + 0.11, 0.36, col)
    text(s, M + 0.4, yy + 0.16, 2.9, 0.3, who, size=14.5, color=INK, bold=True)
    text(s, M + 3.4, yy + 0.17, 4.2, 0.3, can, size=13.5, color=MUTED)
    text(s, M + 7.7, yy + 0.17, 4.1, 0.3, cannot, size=13.5, color=col, bold=True)
    yy += 0.64
footer(s, page=13)
notes(s, """
Access. In most systems you get two roles: editor and reader. That is never enough for a
real institution.

In Dédalo permission is expressed per section AND per field, at three levels: none, read,
or read and write. "None" is genuinely none — the field is not there, and its existence is
not advertised.

That granularity makes three ordinary situations possible. A volunteer catalogues the
physical description of objects and never sees the donor's name. An external researcher
reads an entire collection and cannot alter a character of it. A project team sees its own
material while it stays invisible to the rest of the house until they are ready.

And one thing I will come back to later: these same permissions govern EVERY way into the
data — the staff application, the public service, and any AI assistant you might connect.
There is one set of rules, not one per door.
""")

# ═══ 14 · MEDIA ════════════════════════════════════════════════════════════
s = slide(PAPER, "Files are heritage too")
kicker(s, M, 0.72, "Media")
title(s, "A file is not an attachment.", y=1.12, size=36)
box(s, M, 3.0, 3.5, 1.5, fill=TINT[ORANGE], line=C(0xEE, 0xD6, 0xBE))
text(s, M + 0.4, 3.3, 2.9, 0.4, "THE MASTER", size=11.5, color=ORANGE, bold=True, tracking=150)
text(s, M + 0.4, 3.75, 2.9, 0.5, "Kept as delivered. Never altered.", size=15, color=INK,
     spacing=1.3)
arrow_h(s, M + 3.75, 3.75, 0.75)
for i, o in enumerate(["Derivatives for the web", "Thumbnails for lists",
                       "A transcript, timecoded", "Page references in long texts"]):
    box(s, M + 4.8, 2.6 + i * 0.72, 3.3, 0.6)
    text(s, M + 5.05, 2.78 + i * 0.72, 3.0, 0.3, o, size=13.5, color=INK)
box(s, M + 8.4, 2.6, COL - 8.4, 2.9, fill=WHITE)
accent_bar(s, M + 8.4, 2.88, 2.34, TEAL)
text(s, M + 8.8, 2.86, 2.9, 0.4, "WHY IT MATTERS", size=11.5, color=TEAL, bold=True, tracking=150)
text(s, M + 8.8, 3.32, 2.95, 2.0,
     ["A passage of a recording, with its timecode, becomes a citable unit.", "",
      "Multi-gigabyte video is ordinary: big files are served directly, never carried through the program."],
     size=13, color=MUTED, spacing=1.35, after=2)
footer(s, page=14)
notes(s, """
Files. In most systems a photograph is an attachment — a blob hanging off a record.

Here, an uploaded photograph, recording or video becomes a MASTER, kept exactly as
delivered and never altered, plus all the derivatives the system needs: web versions,
thumbnails, and so on. If your originals are replaced or a new format is needed in ten
years, the master is the thing you go back to.

On top of that, a recording can be transcribed and indexed, so that a PASSAGE — with its
timecode — becomes a citable unit. A researcher does not cite "interview 412"; they cite
the ninety seconds where the thing was actually said. Long texts work the same way with
page references.

And a practical note for whoever manages your storage: large files are served directly by
the web server rather than passed through the program, so a multi-gigabyte video is an
ordinary case here, not an incident.
""")


# ═══ 15 · PART TWO ═════════════════════════════════════════════════════════
s = divider("Part two", "Where it lives, and\nwhat reaches the public",
            "Two separate places, and one decision taken record by record.")
notes(s, """
Part two. Where the system actually lives, and how anything reaches the public.

This is the part your IT department will ask about, and the part that decides how safe you
are. It is simpler than people expect — three slides.
""")

# ═══ 16 · WHERE IT LIVES ═══════════════════════════════════════════════════
s = slide(PAPER, "Where everything lives")
kicker(s, M, 0.72, "Infrastructure, in one picture")
title(s, "Two separate places.\nOnly one faces the internet.", y=1.12)
zw = (COL - 2 * 0.5) / 3
zones = [("INSIDE YOUR INSTITUTION", TEAL, [("Your team", "each with an account"),
                                            ("The Dédalo server", "one program, your network"),
                                            ("The catalogue", "the original — the only one")]),
         ("ON THE WEB SERVER", ORANGE, [("The published copy", "only what you marked"),
                                        ("The public service", "reads, and nothing else"),
                                        ("Files", "CSV · XML · RDF · Markdown")]),
         ("THE WORLD", GREEN, [("Your website", "built from the copy"),
                               ("Visitors, researchers", "whoever you opened it to"),
                               ("AI assistants", "only if you switch them on")])]
for i, (zt, col, rows) in enumerate(zones):
    x = M + i * (zw + 0.5)
    box(s, x, 2.85, zw, 3.35, fill=TINT[col], line=None, radius=0.03)
    text(s, x + 0.35, 3.08, zw - 0.7, 0.3, zt, size=11, color=col, bold=True, tracking=140)
    for j, (h, b) in enumerate(rows):
        yy = 3.55 + j * 0.85
        box(s, x + 0.28, yy, zw - 0.56, 0.72, fill=WHITE, line=None)
        text(s, x + 0.5, yy + 0.12, zw - 1.0, 0.3, h, size=14, color=INK, bold=True)
        text(s, x + 0.5, yy + 0.42, zw - 1.0, 0.25, b, size=11.5, color=MUTED)
    if i < 2:
        arrow_h(s, x + zw + 0.06, 4.5, 0.38)
text(s, M + zw + 0.02, 6.32, 0.5 + zw, 0.3, "you choose what to publish", size=11.5,
     color=MUTED, align=PP_ALIGN.CENTER)
text(s, M + 2 * (zw + 0.5) - 0.5, 6.32, 0.5 + zw, 0.3, "read only — one way", size=11.5,
     color=MUTED, align=PP_ALIGN.CENTER)
footer(s, page=16)
notes(s, """
Here is the whole infrastructure in one picture. Three columns, and the arrows only point
one way.

On the left, inside your institution: your team, each with a personal account; the Dédalo
server, which is one program on one computer on your own network; and the catalogue itself —
every record, every image, every version. That is the original, and the only original there
is. None of it is reachable from the internet.

In the middle, on the web server: a COPY of the records you decided to publish, and the
service that hands that copy out. The same data can also be written as files for a partner
or a national aggregator.

On the right, the world: your website, your visitors, and — only if you switch it on — AI
assistants.

Two things follow. There is no route from the website back into your catalogue, because the
public service cannot write anything anywhere. And the two halves can live on different
machines, so an incident on the public side is not an incident on your archive.
""")

# ═══ 17 · JOURNEY ══════════════════════════════════════════════════════════
s = slide(PAPER, "The journey of one record")
kicker(s, M, 0.72, "From cataloguing to the public")
title(s, "The journey of one record.", y=1.12, size=36)
steps = [("Catalogued", "in the form your\ninstitution\ndesigned", TEAL),
         ("Kept", "with every\nversion, author\nand date", TEAL),
         ("Marked public", "a decision taken\nrecord by record", ORANGE),
         ("Published", "copied outwards,\nin your formats", ORANGE),
         ("Read", "by your website\nand partners", GREEN),
         ("Found by meaning", "if you switch\nthe AI layer on", VIOLET)]
cw = (COL - 5 * 0.22) / 6
for i, (h, b, col) in enumerate(steps):
    x = M + i * (cw + 0.22)
    box(s, x, 2.85, cw, 2.85, fill=WHITE)
    disc(s, x + cw / 2, 3.35, 0.56, col, str(i + 1), size=17)
    text(s, x + 0.14, 3.85, cw - 0.28, 0.6, h, size=15, color=INK, bold=True,
         align=PP_ALIGN.CENTER, spacing=1.12)
    text(s, x + 0.14, 4.72, cw - 0.28, 1.0, b.split("\n"), size=12, color=MUTED,
         align=PP_ALIGN.CENTER, spacing=1.3)
    if i < 5:
        arrow_h(s, x + cw + 0.02, 3.35, 0.18)
box(s, M, 5.95, COL, 0.82, fill=TINT[ORANGE], line=None)
text(s, M + 0.45, 6.12, COL - 1.0, 0.6,
     ["Nothing is published by accident — step 3 is a human decision, and it can be undone.",
      "Nothing is silently lost — step 2 keeps every version, so a mistake is recoverable."],
     size=13, color=INK, spacing=1.4, after=3)
footer(s, page=17)
notes(s, """
Now follow one record — one photograph, one interview, one coin — through its life.

One: someone catalogues it, in the form your institution designed, in as many languages as
you work in. Two: Dédalo keeps it, and keeps every change with its author and date.

Three, and this is the important one: you decide it may be public. Publication is a decision
taken record by record. Until that decision is taken, the record exists only inside the
institution. And un-marking is equally a decision — the next run removes it.

Four: Dédalo publishes. A background job copies the marked records outwards and writes the
formats you asked for. It runs on its own; nobody sits and watches it, and if it is
interrupted it resumes rather than starting over.

Five: your website and your partners read that copy. Six: if you switch the AI layer on, the
same record can be found by what it is about.

Notice what steps 2 and 3 buy you together: you can catalogue frankly — with the donor's
conditions, the internal notes, the unresolved attributions — and still publish confidently.
""")

# ═══ 18 · PUBLISHING IS A DECISION ═════════════════════════════════════════
s = slide(DEEP, "Publishing is a decision")
accent_bar(s, M, 1.85, 2.35, ORANGE, w=0.07)
text(s, M + 0.45, 1.9, COL * 0.85, 1.8,
     ["There is no “make it public” switch", "for the institution."],
     size=38, color=WHITE, bold=True, spacing=1.18)
text(s, M + 0.5, 3.55, COL * 0.72, 0.5,
     "There is a decision per record — and it is reversible.",
     size=21, color=ORANGE_L, bold=True)
cw = (COL - 2 * 0.4) / 3
for i, (h, b) in enumerate([("Catalogue frankly", "donor conditions, internal notes,\nunresolved attributions"),
                            ("Publish selectively", "a record can be public while\nsome of its fields are not"),
                            ("Change your mind", "un-mark it, and the next run\nremoves it from the copy")]):
    x = M + i * (cw + 0.4)
    text(s, x, 4.85, cw, 0.35, h, size=17, color=WHITE, bold=True)
    text(s, x, 5.3, cw, 0.8, b.split("\n"), size=13, color=MUTED_D, spacing=1.35)
notes(s, """
I want to dwell on step three for a moment, because it is where institutions get hurt by
other software.

There is no "make it public" switch for the institution in Dédalo. There is a decision per
record. There is no setting that publishes a whole collection because somebody imported it,
and no default that quietly exposes a field nobody looked at.

That has a liberating effect on cataloguing, and this is the point I most want to land.
Because publication is a separate, deliberate act, your team can catalogue FRANKLY. They can
record the donor's conditions, the internal note about the dubious provenance, the
unresolved attribution, the name of the living person who asked not to be identified. A
record can be public while some of its fields are not.

And you can change your mind. Un-mark it, publish again, it is gone from the public copy.

Catalogue frankly. Publish confidently. Those are not in tension here.
""")

# ═══ 19 · THREE DOORS ══════════════════════════════════════════════════════
s = slide(PAPER, "Three doors")
kicker(s, M, 0.72, "The word “API”, explained once")
title(s, "An API is a door other programs may knock on.", y=1.12)
doors = [("The staff door", "your team, signed in",
          "Everything their permissions\nallow: create, edit, relate,\npublish.",
          "Lets nobody in without\nan account.", TEAL),
         ("The public door", "your website and partners",
          "Reads the published copy:\nrecords, images, search,\npassages.",
          "Writes nothing, anywhere.", ORANGE),
         ("The AI door", "an assistant you connected",
          "Exactly what its user may do:\nsearches, reads, drafts a\nchange.",
          "Saves nothing that nobody\nconfirmed.", VIOLET)]
cw = (COL - 2 * 0.4) / 3
for i, (h, who, can, never, col) in enumerate(doors):
    x = M + i * (cw + 0.4)
    box(s, x, 2.7, cw, 3.45, fill=WHITE)
    box(s, x, 2.7, cw, 0.92, fill=TINT[col], line=None, radius=0.06)
    text(s, x + 0.36, 2.88, cw - 0.7, 0.3, h, size=19, color=INK, bold=True)
    text(s, x + 0.36, 3.22, cw - 0.7, 0.3, who, size=12.5, color=col, bold=True)
    text(s, x + 0.36, 3.85, cw - 0.72, 1.0, can.split("\n"), size=13, color=MUTED, spacing=1.35)
    rule(s, x + 0.36, 5.02, cw - 0.72)
    text(s, x + 0.36, 5.22, cw - 0.72, 0.3, "NEVER", size=10.5, color=ROSE, bold=True, tracking=140)
    text(s, x + 0.36, 5.52, cw - 0.72, 0.6, never.split("\n"), size=13, color=INK, spacing=1.32)
box(s, M, 6.3, COL, 0.5, fill=TINT[TEAL], line=None)
text(s, M + 0.4, 6.44, COL - 0.8, 0.3,
     "One set of rules, checked at every door. Opening the AI door opens nothing else.",
     size=13, color=TEAL, bold=True)
footer(s, page=19)
notes(s, """
You will hear the word "API" in every technical conversation about this, so let me define it
once, in plain terms. An API is a door other programs may knock on. Dédalo has three, and
each is guarded differently.

The staff door: your own team, signed in, doing everything their permissions allow. It lets
nobody in without an account.

The public door: your website, partner portals, researchers you allow. It reads the
published copy — records, images, searches, and passages of texts and recordings. It writes
nothing, anywhere. This door has no handle on the inside.

The AI door: an assistant you connected. It does exactly what the person using it may do —
it searches, it reads, and when it wants to change something it drafts a plan for a human.
It saves nothing that nobody confirmed.

And the line at the bottom is the one to remember when your IT colleagues get nervous:
whoever knocks passes the SAME checks, in the same order, against the same permissions.
Opening the AI door opens nothing else. It is the same lock.
""")


# ═══ 20 · PART THREE ═══════════════════════════════════════════════════════
s = divider("Part three", "Ready for AI —\nand on your terms",
            "Two pieces, already built. Both optional. Both off until you turn them on.",
            color=VIOLET)
notes(s, """
Part three, and this is where I keep the promise I made at the beginning.

Two AI pieces are already built into Dédalo. I want to be precise about what they do and,
just as importantly, what they are structurally unable to do — because in our field that
second list is what decides whether anyone is allowed to use them.

Both are optional. Both are switched off until an administrator turns them on.
""")

# ═══ 21 · THE DAM QUESTION ═════════════════════════════════════════════════
s = slide(PAPER, "The dam question — the callback")
kicker(s, M, 0.72, "Remember the reservoir?", VIOLET)
box(s, M, 1.12, COL, 1.15, fill=WHITE)
accent_bar(s, M, 1.38, 0.63, VIOLET)
text(s, M + 0.5, 1.4, COL - 1.15, 0.8,
     "“Which testimonies talk about people who had to leave their homes because of a dam?”",
     size=19, color=INK, italic=True, spacing=1.35)
lw = (COL - 0.55) / 2
box(s, M, 2.7, lw, 3.5, fill=TINT[SLATE], line=HAIR)
text(s, M + 0.45, 3.0, lw - 0.9, 0.3, "A WORD SEARCH FINDS", size=11.5, color=MUTED,
     bold=True, tracking=150)
text(s, M + 0.45, 3.5, lw - 0.9, 0.5, "Only where somebody wrote the word.", size=16,
     color=MUTED, spacing=1.3)
text(s, M + 0.45, 4.3, lw - 0.9, 0.5, "“…the reservoir works…”", size=17, color=INK, italic=True)
text(s, M + 0.45, 5.35, lw - 0.9, 0.5, "and misses all the rest.", size=15, color=ROSE, bold=True)
box(s, M + lw + 0.55, 2.7, lw, 3.5, fill=TINT[VIOLET], line=C(0xE0, 0xD5, 0xF8), line_w=1.4)
text(s, M + lw + 1.0, 3.0, lw - 0.9, 0.3, "DÉDALO FINDS ALSO", size=11.5, color=VIOLET,
     bold=True, tracking=150)
for i, q in enumerate(["“when the water came and we had to leave”", "“el pantano”",
                       "“they flooded our houses”", "“quan ens van fer marxar”"]):
    text(s, M + lw + 1.0, 3.5 + i * 0.66, lw - 1.0, 0.4, q, size=16, color=INK, italic=True)
footer(s, page=21)
notes(s, """
Remember the reservoir from the second slide? Here it is again.

The researcher asks: which testimonies talk about people who had to leave their homes
because of a dam?

A word search finds only the records where somebody actually wrote the word she typed — "the
reservoir works" — and misses everything else.

Dédalo can also find "when the water came and we had to leave". "El pantano". "They flooded
our houses". "Quan ens van fer marxar". Because those phrases MEAN almost the same thing,
even though they share no words with the question.

This matters more in cultural heritage than almost anywhere else, and I think you all know
why. Our data is multilingual. It is historical — spellings and vocabularies change by
decade and by trade. And it is paraphrastic: humans describe the same object, event or grief
in endlessly different words. An archaeologist's "glazed earthenware vessel with cobalt
decoration" is a curator's "blue-and-white majolica jar" is a donor's "old blue pot".
""")

# ═══ 22 · HOW IT WORKS ═════════════════════════════════════════════════════
s = slide(PAPER, "How semantic search works")
kicker(s, M, 0.72, "Three lines, no mathematics", VIOLET)
title(s, "How it works.", y=1.12, size=36)
cw = (COL - 2 * 0.4) / 3
for i, (t, col) in enumerate([("Dédalo reads the records you\nchoose and stores what they\nmean, as numbers.", VIOLET),
                              ("A question is turned into\nthe same kind of number.", VIOLET),
                              ("The closest meanings come\nback — each citing the record\nit came from.", GREEN)]):
    x = M + i * (cw + 0.4)
    box(s, x, 2.6, cw, 2.35, fill=WHITE)
    disc(s, x + 0.55, 3.1, 0.6, col, str(i + 1), size=18)
    text(s, x + 0.32, 3.68, cw - 0.64, 1.1, t.split("\n"), size=15, color=INK, spacing=1.35)
box(s, M, 5.25, COL, 1.4, fill=TINT[GREEN], line=None)
accent_bar(s, M, 5.5, 0.9, GREEN)
text(s, M + 0.45, 5.46, COL - 1.0, 0.4, "IT ALWAYS CITES", size=11.5, color=GREEN,
     bold=True, tracking=150)
text(s, M + 0.45, 5.88, COL - 1.0, 0.6,
     "Every answer points at the record it came from. Your precise, field-by-field search stays exactly as it is — this is an additional way in, not a replacement.",
     size=14, color=INK, spacing=1.35)
footer(s, page=22)
notes(s, """
How does it do that? Three lines, and no mathematics.

One: Dédalo reads the records you choose — you choose, it is not automatic and not
everything — and stores what they MEAN, as numbers. Two: when someone asks a question, the
question is turned into the same kind of number. Three: the closest meanings come back.

The band at the bottom is what makes this acceptable in a scholarly setting. Every answer
CITES the record it came from. Nothing is summarised into an anonymous paragraph you have to
trust. The researcher opens the record and judges for themselves — which is exactly what
they would have done with a card catalogue, only they found it.

And to be very clear, because this worries people: your precise, structured, field-by-field
search does not go anywhere. It remains the right tool for "every coin minted before 100
BC". Semantic search is an additional way in. In practice people use both in the same
session.
""")

# ═══ 23 · THREE PROMISES ═══════════════════════════════════════════════════
s = slide(PAPER, "AI assistants — three promises")
kicker(s, M, 0.72, "Connecting an assistant", VIOLET)
title(s, "It works with your catalogue —\nit is never given a copy of it.", y=1.12)
cw = (COL - 2 * 0.4) / 3
proms = [("It asks Dédalo", "Never the other way\nround. Your data is not\nhanded to a model to keep."),
         ("It sees what its user sees", "An assistant working for\na colleague sees that\ncolleague's collections —\nno more, ever."),
         ("It proposes; you decide", "It writes a plan you read,\nop by op, and confirm.\nIt cannot save alone.")]
for i, (h, b) in enumerate(proms):
    x = M + i * (cw + 0.4)
    box(s, x, 2.85, cw, 3.35, fill=WHITE)
    accent_bar(s, x, 3.13, 2.79, VIOLET)
    disc(s, x + 0.62, 3.42, 0.56, VIOLET, str(i + 1), size=17)
    text(s, x + 0.36, 4.02, cw - 0.72, 0.75, h, size=18, color=INK, bold=True, spacing=1.12)
    text(s, x + 0.36, 4.9, cw - 0.72, 1.2, b.split("\n"), size=12.5, color=MUTED, spacing=1.35)
footer(s, page=23)
notes(s, """
Now the second piece: connecting an AI assistant — the kind of thing your team may already
be using in a browser tab.

Three promises, and they are structural, not policy.

One: it asks Dédalo, never the other way round. Your archive is not uploaded, not handed to
a model to keep, not used to train anything. The assistant asks a question and gets an
answer, the same way a person does.

Two: it sees only what its user may see. An assistant working for a colleague sees that
colleague's collections, at that colleague's permission level. It is the same lock we saw on
the doors slide.

Three: it proposes, a person decides. When it wants to change something it writes a plan
that you read, operation by operation, and confirm. It cannot save on its own — the part of
the system that talks to the model is built so that it is incapable of writing.

I stress "structural": these are not settings somebody can forget to tick.
""")

# ═══ 24 · SENSITIVE MATERIAL ═══════════════════════════════════════════════
s = slide(DEEP, "Sensitive material stays in the building")
accent_bar(s, M, 1.5, 2.45, ROSE, w=0.07)
text(s, M + 0.45, 1.55, COL * 0.85, 1.8,
     ["Some collections may never", "leave the building."],
     size=40, color=WHITE, bold=True, spacing=1.16)
text(s, M + 0.5, 3.35, COL * 0.66, 0.6,
     "Testimony given under condition. Personal data of living people. Sacred or restricted material.",
     size=15.5, color=MUTED_D, spacing=1.35)
cw = (COL - 2 * 0.4) / 3
for i, (h, b) in enumerate([("Exclude them by name", "They stay fully searchable inside\nthe institution, and reach nothing\noutside it."),
                            ("Or use your own model", "Running on your hardware, so no\nrecord leaves the building at all."),
                            ("One list, both uses", "What you restrict for search is\nrestricted for the assistant too.")]):
    x = M + i * (cw + 0.4)
    box(s, x, 4.45, cw, 2.05, fill=C(0x21, 0x2A, 0x36), line=C(0x33, 0x3E, 0x4C))
    text(s, x + 0.36, 4.7, cw - 0.7, 0.4, h, size=16, color=WHITE, bold=True)
    text(s, x + 0.36, 5.2, cw - 0.7, 1.0, b.split("\n"), size=13, color=MUTED_D, spacing=1.35)
notes(s, """
And now the question every one of you is actually asking, so let us take it head on.

Some collections may never leave the building. Testimony given under condition. Personal
data of living people. Sacred or restricted material. In our field this is not a preference,
it is often a legal or an ethical obligation, and sometimes a promise made to a person who
is still alive.

Dédalo lets you name those collections and exclude them from anything that would reach an
outside model — while they remain fully searchable INSIDE the institution. So your own
researchers keep the benefit; the outside world never touches it.

Or you can point the assistant at a model running on your own hardware, in which case no
record leaves the building at all.

And the two features share one classification: a collection you restrict for search is
restricted for the assistant too. You define it once, and you cannot half-forget it in the
second place.

This is the slide I would put in front of your ethics committee.
""")


# ═══ 25 · RASPA ════════════════════════════════════════════════════════════
s = slide(PAPER, "Measuring what you produce")
kicker(s, M, 0.72, "Data quality")
title(s, "You can argue about quality\nwith evidence.", y=1.12)
text(s, M, 2.55, COL * 0.52, 0.9,
     "Raspa scores a catalogue from 0 to 10 across progressive levels of computational readiness, semantic richness and ethical transparency.",
     size=15, color=MUTED, spacing=1.4)
for i in range(10):
    box(s, M + i * 0.62, 3.85, 0.5, 0.5, fill=(ORANGE if i < 7 else C(0xE4, 0xDE, 0xD4)),
        line=None, radius=0.12)
text(s, M, 4.55, 5.4, 0.4, "structured · modelled · traceable · translatable · open",
     size=12.5, color=MUTED)
for i, (h, b) in enumerate([("In a funding application", "a number behind the argument"),
                            ("Inside the institution", "what to improve next, and why"),
                            ("Between institutions", "a shared, comparable measure")]):
    y = 5.05 + i * 0.6
    box(s, M, y, COL * 0.52, 0.5, fill=WHITE)
    text(s, M + 0.3, y + 0.12, 2.7, 0.3, h, size=13.5, color=INK, bold=True)
    text(s, M + 3.2, y + 0.13, 3.4, 0.3, b, size=13, color=MUTED)
box(s, M + COL * 0.56, 2.5, COL * 0.44, 4.3, fill=TINT[ORANGE], line=None)
text(s, M + COL * 0.56 + 0.45, 2.85, COL * 0.44 - 0.9, 3.5,
     ["“Good data” stops being an opinion.", "",
      "Not because a number captures scholarship — it does not — but because it makes the conversation about quality concrete, repeatable and comparable between projects.",
      "",
      "An extra point is available for data produced with fully free and open tools."],
     size=14.5, color=INK, spacing=1.42, after=4)
footer(s, page=25)
notes(s, """
One smaller thing, but it lands well with funders and with boards.

Dédalo carries a data-quality score called Raspa, which rates a catalogue from zero to ten
across progressive levels: is it structured, is it modelled ontologically, is it traceable,
is it translatable, is it processable with open tools.

I am not going to pretend a number captures scholarship. It does not. What it does is make
the conversation about quality concrete instead of impressionistic. In a funding
application, "our catalogue scores seven, and here is what the eighth point requires" is a
much stronger sentence than "our data is very good". Internally, it tells you what to
improve next and why. And between institutions it gives you something comparable.

There is even an extra point for data produced entirely with free and open tools — which
tells you something about the values of the project.
""")

# ═══ 26 · WHAT YOU OWN ═════════════════════════════════════════════════════
s = slide(PAPER, "What you own")
kicker(s, M, 0.72, "Free and open source, concretely")
title(s, "What that actually means for you.", y=1.12)
owns = [("Your data", "In an open database, on a machine\nyou control. Exportable in full.", TEAL),
        ("Your model", "The shape of your catalogue is yours\n— not a supplier's product decision.", ORANGE),
        ("Your budget", "No licence, no per-seat fee. You pay\nfor hosting and the help you choose.", GREEN),
        ("Your exit", "The one that matters: nothing here\ncan be taken away or discontinued.", VIOLET)]
cw = (COL - 3 * 0.28) / 4
for i, (h, b, col) in enumerate(owns):
    card(s, M + i * (cw + 0.28), 2.8, cw, 2.6, h, b.split("\n"), col, tint=True,
         hsize=18, bsize=12.5)
box(s, M, 5.75, COL, 0.95, fill=TINT[SLATE], line=None)
text(s, M + 0.45, 5.97, COL - 1.0, 0.6,
     "Institutions outlive suppliers. A catalogue built over twenty years should not depend on a company still existing, or on a licence still being renewed.",
     size=15, color=INK, spacing=1.35)
footer(s, page=26)
notes(s, """
Let me translate "free and open source" into things a director cares about, because the
phrase is often heard as "cheap" when what it really means is "yours".

Your data: in an open database, on a machine you control, exportable in full, on any day,
without asking anyone.

Your model: the shape of your catalogue is a decision you take, not a product decision
somebody makes for their whole customer base.

Your budget: no licence and no per-seat fee. What you pay for is hosting, and whatever help
you decide to buy — you can hire anyone, including your own staff.

And your exit, which is the one that really matters: nothing here can be taken away from
you, discontinued, or repriced.

The sentence at the bottom is the argument I would make to a board. Institutions outlive
suppliers. A catalogue built over twenty years should not depend on a company still existing
or a licence still being renewed.
""")

# ═══ 27 · GETTING STARTED ══════════════════════════════════════════════════
s = slide(PAPER, "Getting started")
kicker(s, M, 0.72, "If you wanted to try")
title(s, "What starting actually looks like.", y=1.12)
steps = [("1", "Install it", "On a laptop in ten minutes to look around; on a server, or in a\ncontainer, when you are serious. Nothing to buy first.", TEAL),
         ("2", "Design your model", "The real work, and the part that deserves your best people:\nwhich record types, which fields, which vocabularies.", ORANGE),
         ("3", "Bring your data in", "Import what you already have — spreadsheets, an old database,\ncatalogues in other formats — and start cataloguing.", GREEN),
         ("4", "Decide what is public", "Later, and separately. Publishing is its own decision, and it\ndoes not have to happen in the first year.", VIOLET)]
yy = 2.7
for n, h, b, col in steps:
    box(s, M, yy, COL, 0.95, fill=WHITE)
    accent_bar(s, M, yy + 0.2, 0.55, col)
    disc(s, M + 0.72, yy + 0.47, 0.54, col, n, size=17)
    text(s, M + 1.25, yy + 0.2, 3.3, 0.4, h, size=18, color=INK, bold=True)
    text(s, M + 4.7, yy + 0.19, COL - 5.1, 0.7, b.split("\n"), size=13, color=MUTED,
         spacing=1.3)
    yy += 1.05
footer(s, page=27)
notes(s, """
If any of this appeals, what does starting actually look like? Four steps, and only one of
them is difficult.

One: install it. On a laptop in ten minutes if you just want to click around; on a proper
server or in a container when you are serious. There is nothing to buy first, and no sales
conversation to have.

Two: design your model. This is the real work, and I want to be honest that it IS work.
Which record types, which fields, which vocabularies, which languages. It deserves your best
documentalists, not your IT department — and it is the step where the value of the whole
project is decided.

Three: bring your data in. Spreadsheets, an old database, catalogues in other formats.
Almost nobody starts from nothing.

Four: decide what is public. Later, and separately. Publishing is its own decision and it
does not have to happen in the first year — plenty of institutions catalogue for a long while
before they open anything.
""")

# ═══ 28 · CLOSE ════════════════════════════════════════════════════════════
s = slide(DEEP, "Close")
box(s, 0, 0, W, 0.34, fill=ORANGE, line=None, shape=MSO_SHAPE.RECTANGLE)
text(s, M, 1.95, COL * 0.86, 2.4, ["Catalogue frankly.", "Publish confidently."],
     size=54, color=WHITE, bold=True, spacing=1.16)
rule(s, M, 4.4, 2.2, ORANGE, 0.03)
text(s, M, 4.8, COL * 0.74, 1.4,
     ["Everything else in this talk exists to make those two sentences",
      "possible at the same time — and to keep them true in twenty years."],
     size=18, color=C(0xC8, 0xD1, 0xDC), spacing=1.45)
text(s, M, 6.35, COL * 0.8, 0.4, "Questions — and I am happy to show you the real thing.",
     size=16, color=ORANGE_L, bold=True)
notes(s, """
Let me close where the whole design points.

Catalogue frankly. Publish confidently.

Everything I have shown you exists to make those two sentences true at the same time. The
per-field permissions, so a volunteer can work on an object without seeing the donor. The
per-record publication decision, so an internal note is never an accident waiting to happen.
The version history, so a correction is traceable rather than embarrassing. The separation
of the two servers, so the public side cannot reach back. The links instead of copies, so
the catalogue is still coherent when the people who built it have retired.

And the AI layer, on your terms, so that the knowledge already in your archive can finally
be asked the questions it can answer.

Thank you. I would rather spend the remaining time on your questions than on more slides —
and if it is useful, I am happy to open the real system and show you any of this working.
""")


# ═══ BUILD ═════════════════════════════════════════════════════════════════
pptx_path = os.path.join(OUT_DIR, "dedalo_for_curators.pptx")
prs.save(pptx_path)

words = sum(len(n.split()) for _, n in SCRIPT)
md = ["# Dédalo for curators — speaker script", "",
      "> Generated by `build_deck.py` in this directory — the same source as the speaker",
      "> notes inside `dedalo_for_curators.pptx`. Edit the generator, not this file.", "",
      f"**{len(SCRIPT)} slides · {words} words · about {round(words/130)} minutes** of speaking,",
      "so roughly 30 with pauses. The three dividers are quick; give the time to slides 7,",
      "10, 18 and 24. If you are short, slides 25 and 26 can be dropped without breaking",
      "the argument.", ""]
for i, (t, n) in enumerate(SCRIPT, 1):
    md += [f"## {i}. {t}", "", n, ""]
open(os.path.join(OUT_DIR, "speaker_script.md"), "w").write("\n".join(md))

print(f"{len(SCRIPT)} slides · {words} script words (~{words/130:.0f} min)")
print("wrote", pptx_path)
