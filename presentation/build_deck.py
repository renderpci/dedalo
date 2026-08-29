#!/usr/bin/env python3
"""
Dédalo — the curator's presentation.

ONE source for both artifacts: this file builds `dedalo_for_curators.pptx`
(slides + embedded speaker notes) and `speaker_script.md` (the same notes,
printable). Edit here and re-run; never edit the .pptx by hand and expect the
script to follow.

    <venv>/bin/python presentation/build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
W, H = 13.333, 7.5
M = 0.9                      # page margin
COL = W - 2 * M              # content width

# ── palette ────────────────────────────────────────────────────────────────
INK    = C(0x1A, 0x1C, 0x22)
PAPER  = C(0xFA, 0xF8, 0xF4)
DEEP   = C(0x14, 0x1A, 0x24)
DEEP2  = C(0x1E, 0x28, 0x36)
ORANGE = C(0xC2, 0x62, 0x2A)
ORANGE_L = C(0xE8, 0x9A, 0x3C)
TEAL   = C(0x0E, 0x74, 0x90)
GREEN  = C(0x15, 0x80, 0x3D)
VIOLET = C(0x6D, 0x28, 0xD9)
ROSE   = C(0xBE, 0x12, 0x3C)
MUTED  = C(0x66, 0x6E, 0x7A)
MUTED_D= C(0x93, 0x9E, 0xAC)
HAIR   = C(0xE2, 0xDC, 0xD1)
WHITE  = C(0xFF, 0xFF, 0xFF)
TINT   = {ORANGE: C(0xFB, 0xF1, 0xE7), TEAL: C(0xEC, 0xF6, 0xF9),
          GREEN: C(0xED, 0xF8, 0xF1), VIOLET: C(0xF4, 0xEF, 0xFD),
          ROSE:  C(0xFD, 0xEF, 0xF2), INK: C(0xF2, 0xF0, 0xEC)}

FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]
SCRIPT = []          # (title, notes) collected for the printable script


# ── primitives ─────────────────────────────────────────────────────────────
def slide(bg=PAPER, title_for_script=""):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    s._script_title = title_for_script
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()
    SCRIPT.append((getattr(s, "_script_title", ""), text.strip()))


def _spc(run, hundredths_pt):
    """Letter-spacing — python-pptx has no API for it, so set the attribute."""
    run.font._element.set("spc", str(hundredths_pt))   # rPr/@spc, in 1/100 pt


def text(s, x, y, w, h, lines, size=14, color=MUTED, bold=False, align=PP_ALIGN.LEFT,
         spacing=1.25, after=0, tracking=None, italic=False, font=FONT):
    """lines: a string, or a list of strings / (string, overrides-dict) pairs."""
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        over = {}
        if isinstance(ln, tuple):
            ln, over = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.line_spacing = over.get("spacing", spacing)
        p.space_after = Pt(over.get("after", after))
        r = p.add_run()
        r.text = ln
        r.font.name = over.get("font", font)
        r.font.size = Pt(over.get("size", size))
        r.font.bold = over.get("bold", bold)
        r.font.italic = over.get("italic", italic)
        r.font.color.rgb = over.get("color", color)
        tr = over.get("tracking", tracking)
        if tr:
            _spc(r, tr)
    return tb


def est_lines(content, w_in, size):
    """Rough wrapped-line count. Helvetica Neue averages ~0.50 em per character,
    so chars-per-line = width(pt) / (0.50 * size). Renderers differ slightly, so
    this deliberately errs on the generous side."""
    cpl = max(8, int(w_in * 72 / (size * 0.52)))
    items = content if isinstance(content, list) else [content]
    n = 0
    for it in items:
        t = it[0] if isinstance(it, tuple) else it
        n += max(1, -(-len(t) // cpl))
    return n


def line_h(size, spacing=1.3):
    return size * spacing / 72.0


def kicker(s, x, y, label, color=ORANGE, size=11.5):
    return text(s, x, y, COL, 0.3, label.upper(), size=size, color=color,
                bold=True, tracking=160)


def title(s, txt, x=M, y=0.95, w=None, size=36, color=INK, spacing=1.05):
    return text(s, x, y, w or COL, 1.6, txt, size=size, color=color, bold=True,
                spacing=spacing)


def rule(s, x, y, w, color=HAIR, thickness=0.014):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y), In(w), In(thickness))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def box(s, x, y, w, h, fill=WHITE, line=HAIR, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.045, line_w=1.0):
    sh = s.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    sh.text_frame.text = ""
    return sh


def accent_bar(s, x, y, h, color, w=0.055):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.5
    except Exception:
        pass
    return sh


def disc(s, cx, cy, d, fill, label="", color=WHITE, size=18, bold=True):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, In(cx - d / 2), In(cy - d / 2), In(d), In(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    if label:
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return sh


def shape_icon(s, kind, cx, cy, size, color):
    """A pictogram from a PowerPoint preset — editable, never a raster."""
    presets = {
        "db": MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
        "doc": MSO_SHAPE.FLOWCHART_DOCUMENT,
        "docs": MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
        "cloud": MSO_SHAPE.CLOUD,
        "screen": MSO_SHAPE.FLOWCHART_DISPLAY,
        "cube": MSO_SHAPE.CUBE,
        "bubble": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        "diamond": MSO_SHAPE.DIAMOND,
    }
    sh = s.shapes.add_shape(presets[kind], In(cx - size / 2), In(cy - size / 2),
                            In(size), In(size * 0.82))
    sh.fill.background()
    sh.line.color.rgb = color; sh.line.width = Pt(1.9)
    sh.shadow.inherit = False
    from pptx.oxml.ns import qn
    st = sh._element.find(qn('p:style'))
    if st is not None and st.find(qn('a:effectRef')) is not None:
        st.remove(st.find(qn('a:effectRef')))
    return sh


def arrow_h(s, x, y, length, color=C(0x9A, 0xA4, 0xB2), thick=0.028, head=0.16):
    stem = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y - thick / 2),
                              In(length - head * 0.75), In(thick))
    stem.fill.solid(); stem.fill.fore_color.rgb = color
    stem.line.fill.background(); stem.shadow.inherit = False
    tip = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                             In(x + length - head), In(y - head / 2), In(head), In(head))
    tip.fill.solid(); tip.fill.fore_color.rgb = color
    tip.line.fill.background(); tip.shadow.inherit = False
    tip.rotation = 90
    return stem


def card(s, x, y, w, h, heading, lines, color=ORANGE, tint=False, num=None,
         hsize=18, bsize=13):
    box(s, x, y, w, h, fill=(TINT[color] if tint else WHITE))
    accent_bar(s, x + 0.0, y + 0.28, h - 0.56, color)
    tx = x + 0.34
    ty = y + 0.3
    if num is not None:
        disc(s, tx + 0.24, ty + 0.22, 0.48, color, str(num), size=16)
        ty += 0.72
    hl = est_lines(heading, w - 0.62, hsize)
    text(s, tx, ty, w - 0.62, hl * line_h(hsize, 1.15) + 0.1, heading, size=hsize,
         color=INK, bold=True, spacing=1.12)
    if lines:
        text(s, tx, ty + hl * line_h(hsize, 1.15) + 0.16, w - 0.62, h - 1.0, lines,
             size=bsize, color=MUTED, spacing=1.32, after=3)


def footer(s, txt="Dédalo · Cultural Heritage management", color=None, page=None):
    text(s, M, H - 0.62, COL * 0.7, 0.3, txt, size=10, color=color or C(0xA6, 0x9E, 0x92),
         tracking=60)
    if page is not None:
        text(s, W - M - 1.2, H - 0.62, 1.2, 0.3, str(page), size=10,
             color=color or C(0xA6, 0x9E, 0x92), align=PP_ALIGN.RIGHT)


def line(s, x1, y1, x2, y2, color=HAIR, w=1.2):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(x1), In(y1), In(x2), In(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    ln.shadow.inherit = False
    return ln


def bar(s, x, y, w, h, fill=HAIR):
    return box(s, x, y, w, h, fill=fill, line=None, radius=min(0.5, h / 2))


def chip(s, x, y, w, h, label, fill=WHITE, line_c=HAIR, color=INK, size=12, bold=False):
    sh = box(s, x, y, w, h, fill=fill, line=line_c, radius=0.5)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def ring(s, cx, cy, d, line_c, label="", color=INK, size=12, fill=WHITE, line_w=1.7,
         bold=True):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, In(cx - d / 2), In(cy - d / 2), In(d), In(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line_c
    sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if label:
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return sh


def ctext(s, y, txt, size=15, color=INK, bold=False, italic=False, tracking=None):
    return text(s, 0, y, W, 0.5, txt, size=size, color=color, bold=bold, italic=italic,
                align=PP_ALIGN.CENTER, tracking=tracking)


def dots(s, cx, cy, d, color):
    return disc(s, cx, cy, d, color)


# ═════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═════════════════════════════════════════════════════════════════════════

# ── 1 · Title ──────────────────────────────────────────────────────────────
s = slide(DEEP, "Dédalo — the living catalogue of your collection")
LINK_D = C(0x39, 0x47, 0x5A)
for (x1, y1, x2, y2) in [(9.3, 1.7, 11.4, 2.4), (11.4, 2.4, 10.3, 3.6),
                          (10.3, 3.6, 12.1, 4.4), (10.3, 3.6, 9.6, 5.2),
                          (9.6, 5.2, 11.2, 6.0), (11.2, 6.0, 12.1, 4.4),
                          (9.3, 1.7, 10.3, 3.6)]:
    line(s, x1, y1, x2, y2, color=LINK_D, w=1.2)
disc(s, 9.3, 1.7, 0.5, ORANGE)
disc(s, 11.4, 2.4, 0.34, TEAL)
disc(s, 10.3, 3.6, 0.64, ORANGE_L)
disc(s, 12.1, 4.4, 0.3, VIOLET)
disc(s, 9.6, 5.2, 0.44, TEAL)
disc(s, 11.2, 6.0, 0.26, ORANGE)
for (x, y) in [(9.95, 2.75), (11.75, 3.15), (10.55, 4.95), (12.3, 5.5)]:
    disc(s, x, y, 0.1, MUTED_D)
text(s, M, 1.62, 8.0, 0.4, "CULTURAL HERITAGE · ORAL MEMORY · OPEN SOURCE",
     size=13, color=ORANGE_L, bold=True, tracking=220)
text(s, M - 0.06, 2.05, 8.5, 1.8, "Dédalo", size=96, color=PAPER, bold=True, spacing=1.0)
text(s, M, 3.85, 7.6, 0.6, "The living catalogue of your collection.",
     size=25, color=MUTED_D, spacing=1.1)
rule(s, M, 4.75, 2.3, color=ORANGE_L, thickness=0.03)
text(s, M, 5.0, 7.5, 0.4,
     "A guided tour for curators, archivists and keepers of memory — about 30 minutes, no code.",
     size=14, color=MUTED_D, spacing=1.2)
footer(s, "Dédalo · dedalo.dev", color=C(0x5C, 0x6A, 0x7C))
notes(s, """
Welcome, and thank you for the half hour you're giving us.

Dédalo is a system for managing cultural heritage — collections, inventories,
archives, oral testimony. Coins and amphorae, interviews and photographs,
places and people.

I promise you two things today. First: no code. Not one line. Everything I
show will be about your work — cataloguing, researching, publishing — and not
about servers. Second: by the end, you'll be able to explain what Dédalo is,
and why it is different from the catalogue software you've met before, in
about three sentences.

The name is Dédalo — Daedalus, the craftsman. Keep that image in mind: not a
box that stores things, but a workshop that shapes them.

Let's begin with an uncomfortable question: why does the heritage field keep
producing new software, when there is already so much of it?
""")

# ── 2 · Two ways to fail ───────────────────────────────────────────────────
s = slide(PAPER, "Heritage software usually fails in one of two ways")
kicker(s, M, 0.62, "WHY ANOTHER SYSTEM")
title(s, "Software usually fails in one of two ways.", size=36)
# left card — the pretty spreadsheet
box(s, M, 2.05, 5.55, 3.85, WHITE)
accent_bar(s, M, 2.33, 3.29, ROSE)
text(s, M + 0.35, 2.32, 5.0, 0.4, "The pretty spreadsheet", size=18, color=INK, bold=True)
text(s, M + 0.35, 2.78, 5.0, 0.4, "fast to start — incoherent after three years",
     size=12, color=MUTED)
for r in range(3):
    for c in range(5):
        x = M + 0.55 + c * 0.62
        y = 3.55 + r * 0.66
        if (r, c) == (1, 2):
            continue
        sh = box(s, x, y, 0.44, 0.44, fill=(HAIR if (r + c) % 2 else None),
                 line=(None if (r + c) % 2 else MUTED_D), shape=MSO_SHAPE.RECTANGLE)
        if (r, c) == (0, 3):
            sh.rotation = 9
        if (r, c) == (2, 1):
            sh.rotation = -7
text(s, M + 2.25, 4.45, 0.6, 0.5, "?", size=26, color=ROSE, bold=True,
     align=PP_ALIGN.CENTER)
# right card — the rigid database
box(s, 6.88, 2.05, 5.55, 3.85, WHITE)
accent_bar(s, 6.88, 2.33, 3.29, INK)
text(s, 7.23, 2.32, 5.0, 0.4, "The rigid database", size=18, color=INK, bold=True)
text(s, 7.23, 2.78, 5.1, 0.4,
     "designed by someone who never catalogued anything — adding a field is a project",
     size=12, color=MUTED)
for r in range(3):
    for c in range(5):
        box(s, 7.43 + c * 0.62, 3.55 + r * 0.66, 0.44, 0.44,
            fill=C(0xE6, 0xE2, 0xDA), line=C(0xD4, 0xCE, 0xC2), shape=MSO_SHAPE.RECTANGLE)
shk = s.shapes.add_shape(MSO_SHAPE.OVAL, In(9.64), In(3.6), In(0.44), In(0.44))
shk.fill.background()
shk.line.color.rgb = INK
shk.line.width = Pt(3.2)
shk.shadow.inherit = False
box(s, 9.55, 3.92, 0.62, 0.52, fill=INK, line=None, radius=0.22)
ctext(s, 6.35, "Dédalo was built to avoid both.", size=16, color=ORANGE, bold=True)
footer(s, page=2)
notes(s, """
Software for cultural heritage tends to fail in one of two ways, and most of
us have been burned by both.

The first failure is the pretty spreadsheet. It's wonderful on day one.
Everyone adds columns differently, every volunteer types dates their own way.
After three years nobody is sure which file is the real one, and the person
who understands it has left the institution. Fast to start — impossible to
keep coherent.

The second failure is the rigid database. Someone who has never held an
amphora designed its tables. It is well behaved, and completely inflexible:
you need one extra field for a new project, and that small human request
becomes a ticket, a developer, a budget, a release, a delay. A project to
change something that should take a minute.

Notice what both failures have in common: the shape of your catalogue is not
in your hands. Either it's in nobody's hands, or it's in a programmer's
hands. Hold that thought — because Dédalo's central idea is precisely to give
that shape back to you.
""")

# ── 3 · Three purposes ─────────────────────────────────────────────────────
s = slide(PAPER, "It exists to do three things")
kicker(s, M, 0.62, "WHAT IS DÉDALO")
title(s, "It exists to do three things.", size=36)
pillars = [
    ("01", "doc", ORANGE, "Produce good data",
     "structured, linked, in every language"),
    ("02", "db", TEAL, "Keep it for decades",
     "every version, every author, every date"),
    ("03", "cloud", VIOLET, "Publish what should be public",
     "a decision — never an accident"),
]
for i, (num, icon, col, head, sub) in enumerate(pillars):
    x = M + i * 3.99
    cx = x + 1.83
    text(s, x, 2.25, 3.7, 0.35, num, size=13, color=MUTED_D, bold=True,
         align=PP_ALIGN.CENTER, tracking=200)
    shape_icon(s, icon, cx, 3.15, 1.05, col)
    text(s, x, 3.95, 3.7, 0.75, head, size=19, color=INK, bold=True,
         align=PP_ALIGN.CENTER, spacing=1.05)
    text(s, x, 4.75, 3.7, 0.6, sub, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    rule(s, cx - 0.45, 5.45, 0.9, color=col, thickness=0.025)
ctext(s, 6.15, "Not a website builder. Not a generic CMS. A tool to produce good data",
      size=13.5, color=MUTED)
ctext(s, 6.45, "and keep it trustworthy for thirty years.", size=13.5, color=MUTED)
footer(s, page=3)
notes(s, """
So what is Dédalo? Strip away everything else and it has exactly three jobs.

One — produce good data. Not "data", good data: structured, consistent,
linked to each other, written in every language your institution works in.
This is the quiet, unglamorous heart of the whole project.

Two — keep it for decades. A collection outlives spreadsheets, outlives
staff, outlives software vendors. Dédalo keeps every version of every record
with its author and its date, so the catalogue is still trustworthy — and
still explainable — in thirty years.

Three — publish what should be public. Cataloguing and publishing are two
different acts, and Dédalo never confuses them. What goes out to the world is
a deliberate decision, taken record by record.

Notice what is NOT on this list: Dédalo is not a website builder, and it's
not a generic content manager. Your website will talk to Dédalo — but
Dédalo's job is the collection itself: its quality, its memory, its safety.
""")

# ── 4 · Two rooms ──────────────────────────────────────────────────────────
s = slide(PAPER, "One system, two rooms")
kicker(s, M, 0.62, "THE SHAPE OF IT")
title(s, "One system. Two rooms.", size=36)
box(s, M, 2.05, 5.3, 3.75, DEEP2)
text(s, M + 0.35, 2.35, 4.6, 0.35, "THE WORKSHOP", size=12, color=ORANGE_L, bold=True,
     tracking=200)
text(s, M + 0.35, 2.75, 4.7, 0.5, "where the collection is made", size=19, color=PAPER,
     bold=True)
text(s, M + 0.35, 3.55, 4.7, 2.0,
     ["Your team, each with a personal account",
      "Your own network — not the internet",
      "The original. The only original there is."],
     size=12.5, color=MUTED_D, spacing=1.15, after=10)
box(s, 7.13, 2.05, 5.3, 3.75, TINT[TEAL], line=TEAL)
text(s, 7.48, 2.35, 4.6, 0.35, "THE GALLERY", size=12, color=TEAL, bold=True,
     tracking=200)
text(s, 7.48, 2.75, 4.7, 0.5, "where the public visits", size=19, color=INK, bold=True)
text(s, 7.48, 3.55, 4.7, 2.0,
     ["A published copy — read-only",
      "Your website, partners, researchers, AI",
      "Shares nothing with the workshop"],
     size=12.5, color=MUTED, spacing=1.15, after=10)
arrow_h(s, 6.28, 3.98, 0.86, color=ORANGE, thick=0.04, head=0.24)
text(s, 5.9, 3.45, 1.6, 0.3, "publish", size=12, color=ORANGE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, 5.9, 4.22, 1.6, 0.3, "one way only", size=10, color=MUTED,
     align=PP_ALIGN.CENTER)
ctext(s, 6.3, "The door to the gallery has no handle on the inside.", size=15.5,
      color=INK, italic=True)
footer(s, page=4)
notes(s, """
Here is the shape of the whole thing, and the museum analogy is the project's
own. Dédalo is two rooms.

Behind, the workshop. Your team works there on your own network — you do not
even need the internet to catalogue. It holds the catalogue itself: every
record, every image, every version. This is the original, and the only
original there is. It never faces the internet.

In front, the gallery. That's your public website, partners, researchers —
even AI assistants. What they see is a published copy: flattened, read-only,
sitting on a different machine.

Data flows one way: workshop to gallery. There is no route back in. The
public door has no handle on the inside — it can hand things out, but it
structurally cannot write anything, anywhere.

This one design decision answers most of the scary questions. If the public
site is attacked or goes down, cataloguing continues untouched. If your
network is down, the website keeps serving. And nothing — nothing — is ever
public by accident.
""")

# ── 5 · Active ontology ────────────────────────────────────────────────────
s = slide(PAPER, "The shape of your catalogue is data")
kicker(s, M, 0.62, "THE IDEA THAT DRIVES IT")
title(s, "The shape of your catalogue is data.", size=36)
box(s, M, 2.05, 5.45, 3.7, WHITE)
accent_bar(s, M, 2.33, 3.14, MUTED_D)
text(s, M + 0.35, 2.3, 4.8, 0.4, "Everywhere else", size=17, color=INK, bold=True)
chain = ["“We need one new field”", "a developer", "a release", "a budget", "a delay"]
for i, t in enumerate(chain):
    chip(s, M + 1.35, 2.85 + i * 0.58, 3.0, 0.44, t,
         fill=(TINT[INK] if i == 0 else WHITE),
         line_c=(HAIR if i == 0 else C(0xD8, 0xD3, 0xC8)),
         color=(MUTED if i else INK), size=11.5)
    if i < 4:
        text(s, M + 2.6, 3.18 + i * 0.58, 0.5, 0.3, "↓", size=13, color=MUTED_D,
             align=PP_ALIGN.CENTER)
box(s, 7.0, 2.05, 5.45, 3.7, TINT[ORANGE])
accent_bar(s, 7.0, 2.33, 3.14, ORANGE)
text(s, 7.35, 2.3, 4.8, 0.4, "In Dédalo", size=17, color=INK, bold=True)
chain2 = ["“We need one new field”", "you edit the ontology", "the form changes — live"]
for i, t in enumerate(chain2):
    last = i == 2
    chip(s, 8.0 + 0.0, 2.95 + i * 0.85, 3.5, 0.52, t,
         fill=(ORANGE if last else WHITE),
         line_c=(ORANGE if last else HAIR),
         color=(WHITE if last else INK), size=12.5, bold=last)
    if i < 2:
        text(s, 9.5, 3.42 + i * 0.85, 0.5, 0.3, "↓", size=13, color=ORANGE,
             align=PP_ALIGN.CENTER)
ctext(s, 6.3, "Adding a field is an act of documentation — not of programming.",
      size=15.5, color=INK, bold=True)
footer(s, page=5)
notes(s, """
Now the central idea. If you remember one thing from today, make it this one.

In almost every system, the shape of the catalogue — which kinds of records
exist, which fields they have, in what order, with which vocabularies —
lives inside the program. Change the shape, and you change the program:
developer, release, budget, delay.

In Dédalo, the shape of the catalogue is itself data. It lives in something
called the ontology — think of it as the catalogue's constitution — and you
edit it inside Dédalo, like anything else. Adding a field, creating a whole
new kind of record, making a text box into a controlled vocabulary,
reordering a form: these become acts of documentation, not programming.

Why does this matter to an institution? Because your descriptive model can
follow your discipline instead of the software's assumptions — and it can
evolve. A modest inventory can grow into a full scientific catalogue without
rebuilding anything. A standard changes? You update the model yourself.

One honest caveat: the power moves to whoever designs your model. The design
deserves the same care as a cataloguing manual. Dédalo stops the software
imposing a bad model on you — it can't stop you designing one.
""")

# ── 6 · A record ───────────────────────────────────────────────────────────
s = slide(PAPER, "Every object gets a complete page")
kicker(s, M, 0.62, "A DAY IN THE WORKSHOP")
title(s, "Every object gets a complete page.", size=36)
box(s, 1.05, 2.1, 7.0, 4.15, WHITE, line=C(0xD4, 0xCE, 0xC2))
hd = box(s, 1.05, 2.1, 7.0, 0.56, INK, line=None, radius=0.09)
text(s, 1.35, 2.22, 4.0, 0.35, "Coins · inv. 1234", size=13, color=PAPER, bold=True)
text(s, 6.2, 2.22, 1.7, 0.35, "CA  ES  EN", size=11, color=ORANGE_L, bold=True,
     align=PP_ALIGN.RIGHT)
rows = [
    (3.0, "Title", [("bar", 4.3, HAIR)]),
    (3.62, "Date", [("bar", 1.9, HAIR)]),
    (4.24, "Material", [("term", "bronze"), ("term", "silver")]),
    (4.92, "Photography", [("img",)]),
    (5.85, "Found in", [("term2", "→  Tossal Rodó · site")]),
]
for y, label, elems in rows:
    chip(s, 1.35, y, 1.35, 0.36, label, fill=TINT[INK], line_c=None, color=INK,
         size=10.5, bold=True)
    x = 3.0
    for e in elems:
        if e[0] == "bar":
            bar(s, x, y + 0.07, e[1], 0.22)
            x += e[1] + 0.2
        elif e[0] == "term":
            chip(s, x, y, 1.15, 0.36, e[1], fill=TINT[TEAL], line_c=None, color=TEAL,
                 size=10.5, bold=True)
            x += 1.35
        elif e[0] == "term2":
            chip(s, x, y, 2.6, 0.36, e[1], fill=TINT[VIOLET], line_c=None,
                 color=VIOLET, size=10.5, bold=True)
        elif e[0] == "img":
            ib = box(s, x, y - 0.28, 1.15, 0.92, TINT[TEAL], line=TEAL)
            disc(s, x + 0.28, y - 0.06, 0.16, ORANGE_L)
            tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, In(x + 0.36),
                                     In(y + 0.13), In(0.62), In(0.46))
            tri.fill.solid()
            tri.fill.fore_color.rgb = TEAL
            tri.line.fill.background()
            tri.shadow.inherit = False
cards = [
    ("Sections", "kinds of record: Coins, Testimonies, Places, People", ORANGE),
    ("Components", "fields that know what they are — dates, places, media, links",
     TEAL),
    ("One home", "the object, its images, its history — all on one page", VIOLET),
]
for i, (head, sub, col) in enumerate(cards):
    card(s, 8.55, 2.1 + i * 1.45, 3.85, 1.28, head, [sub], color=col, hsize=15.5,
         bsize=11.5)
ctext(s, 6.5, "The form is the one your institution designed — not the one software imposed.",
      size=13.5, color=MUTED, italic=True)
footer(s, page=6)
notes(s, """
Let's stand in the workshop for a minute and look over a cataloguer's
shoulder.

This is a record — one coin. Above it, its section: the kind of record it is.
Coins is a section. Oral testimonies is a section. People, places, excavation
sites — sections. If you're an archivist, think "series"; if you've met
databases, think "table", but a table that can be re-shaped by you.

Inside the record are components — fields, but fields that understand their
kind. The date component understands imperfect dates — "before 100 BC,
probably". The geographic component understands points and areas. The media
component knows a file has a master, derivatives and a thumbnail. The
relation component knows its content is a pointer to another record, not a
piece of text.

You'll notice tabs at the top: CA, ES, EN — we'll come to languages shortly.

And remember the last slide: every element of this form — which fields exist,
their order, their vocabularies — was designed by the institution, not
imposed by a programmer. This page is yours.
""")

# ── 7 · Links, not copies ──────────────────────────────────────────────────
s = slide(PAPER, "Correct once, right everywhere")
kicker(s, M, 0.62, "LINKS, NOT COPIES")
title(s, "Correct once. Right everywhere.", size=36)
HUB = (3.9, 4.05)
sats = [(6.1, 3.15), (6.5, 4.75), (5.35, 5.75), (2.75, 5.85), (1.55, 4.85),
        (1.55, 3.15), (2.75, 2.25), (5.1, 2.25)]
for (x, y) in sats:
    line(s, HUB[0], HUB[1], x, y, color=C(0xC9, 0xC2, 0xB4), w=1.1)
for (x, y) in sats:
    box(s, x - 0.16, y - 0.16, 0.32, 0.32, fill=WHITE, line=MUTED_D)
disc(s, HUB[0], HUB[1], 1.55, ORANGE, "Vall de Cabó", size=12)
text(s, 2.35, 5.05, 3.1, 0.3, "one term · one record", size=11, color=ORANGE,
     bold=True, align=PP_ALIGN.CENTER)
text(s, 2.42, 5.32, 3.0, 0.3, "900 records point at it", size=11, color=MUTED,
     align=PP_ALIGN.CENTER)
ideas = [
    ("Fix it once", "change the spelling, the date, the name — every record that points here is right, immediately", ORANGE),
    ("Ask backwards", "“which testimonies mention this person?” “what came from this excavation?” — questions you never planned", TEAL),
    ("No drift", "two records can't hold two spellings of the same place, because they hold no spelling at all", VIOLET),
]
for i, (head, sub, col) in enumerate(ideas):
    card(s, 7.65, 2.1 + i * 1.5, 4.75, 1.35, head, [sub], color=col, hsize=16,
         bsize=11)
ctext(s, 6.5, "This is what keeps a catalogue usable after twenty years.", size=14,
      color=INK, italic=True)
footer(s, page=7)
notes(s, """
This slide is the single biggest difference between a catalogue and a
spreadsheet.

When a record mentions a place in a spreadsheet, you type the string "Vall de
Cabó". Nine hundred records, nine hundred copies of the name. One typo, one
reform of the toponym, one clerical error — and you have nine hundred little
problems.

Dédalo never copies names. It stores a pointer. "Vall de Cabó" is not a
string repeated everywhere — it is one record, with its own description,
variants and history, and nine hundred records point at it.

Three consequences, and they're the ones every documentalist has wanted
forever. Correct once: fix the term, and every record is right, immediately,
everywhere. Ask backwards: because the link is real, you can ask "which
testimonies mention this person?" or "what came from this excavation?" —
questions nobody planned in advance. No drift: two records cannot disagree
about a spelling they don't hold.

Links, not copies. It sounds small. It is the whole difference between a list
and a network of knowledge.
""")

# ── 8 · Thesauri ───────────────────────────────────────────────────────────
s = slide(PAPER, "Words that mean something")
kicker(s, M, 0.62, "CONTROLLED VOCABULARIES")
title(s, "Words that mean something.", size=36)
TC = C(0xA9, 0xCC, 0xD6)
root = (4.1, 2.62)
kids = [(2.15, 4.2), (4.1, 4.2), (6.05, 4.2)]
gk = [(3.35, 5.5), (4.85, 5.5)]
for k in kids:
    line(s, root[0], root[1], k[0], k[1], color=TC, w=1.4)
for g in gk:
    line(s, kids[1][0], kids[1][1], g[0], g[1], color=TC, w=1.4)
ring(s, root[0], root[1], 1.25, TEAL, "Materials", size=12)
for lbl, (x, y) in zip(["Stone", "Metal", "Ceramic"], kids):
    ring(s, x, y, 1.1, TEAL, lbl, size=12)
for lbl, (x, y) in zip(["Bronze", "Silver"], gk):
    ring(s, x, y, 0.95, TEAL, lbl, size=11)
text(s, 1.0, 6.12, 6.2, 0.3, "a thesaurus: a tree of terms, one branch per domain",
     size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
box(s, 7.65, 2.1, 4.75, 3.9, WHITE)
accent_bar(s, 7.65, 2.38, 3.34, TEAL)
text(s, 8.0, 2.35, 4.2, 0.45, "A term is a record.", size=19, color=INK, bold=True)
text(s, 8.0, 2.95, 4.15, 1.3,
     "It can be described, dated, related, translated and cited — exactly like an object.",
     size=13, color=MUTED, spacing=1.3)
text(s, 8.0, 4.15, 4.15, 0.3, "VOCABULARIES YOU CAN BUILD", size=10, color=MUTED_D,
     bold=True, tracking=150)
tx = 8.0
for lbl in ["places", "people"]:
    chip(s, tx, 4.5, 1.25, 0.4, lbl, fill=TINT[TEAL], line_c=None, color=TEAL, size=11,
         bold=True)
    tx += 1.45
tx = 8.0
for lbl in ["subjects", "materials", "techniques"]:
    chip(s, tx, 5.0, 1.35, 0.4, lbl, fill=TINT[TEAL], line_c=None, color=TEAL, size=11,
         bold=True)
    tx += 1.5
ctext(s, 6.5, "One agreed language — for the whole institution, for the whole web of data.",
      size=14, color=INK, italic=True)
footer(s, page=8)
notes(s, """
Links need something to link to — and that brings us to vocabularies.

A thesaurus in Dédalo is a tree of terms: places, personal names, subjects,
materials, techniques, typologies. Controlled vocabulary means: when nine
cataloguers write "the Iron Age", they all reach for the same node — not nine
slightly different strings.

The important part is subtle: the terms in the tree are ordinary records.
Remember the previous slide — a pointer, not a copy? The term record itself
can be described, dated, translated, cited. Your toponymic thesaurus entry for
a valley can carry its coordinates, its alternative names through history,
its translations. It is documented as carefully as any museum object,
because in Dédalo it is one.

And because vocabularies are yours — remember, the shape of the catalogue is
data — when a new controlled vocabulary appears in your discipline, or a
national aggregator imposes one, adopting it is a documentation task you
perform, not a feature request you file.

Words that mean something, agreed once, used everywhere.
""")

# ── 9 · Languages ──────────────────────────────────────────────────────────
s = slide(PAPER, "Every value, in every language")
kicker(s, M, 0.62, "MANY LANGUAGES")
title(s, "Every value, in every language.", size=36)
box(s, 1.05, 2.1, 6.3, 3.95, WHITE, line=C(0xD4, 0xCE, 0xC2))
tabs = [("CA", True), ("ES", False), ("EN", False), ("FR", False)]
for i, (t, on) in enumerate(tabs):
    chip(s, 1.4 + i * 0.85, 2.35, 0.7, 0.4, t, fill=(ORANGE if on else WHITE),
         line_c=(ORANGE if on else HAIR), color=(WHITE if on else MUTED), size=11.5,
         bold=True)
text(s, 5.0, 2.42, 2.1, 0.3, "· the record's data", size=11, color=MUTED)
text(s, 1.4, 3.05, 2.0, 0.3, "TITLE", size=9.5, color=MUTED_D, bold=True, tracking=150)
langs = [("CA", "Àmfora grega"), ("ES", "Ánfora griega"), ("EN", "Greek amphora"),
         ("FR", "Amphore grecque")]
for i, (lg, val) in enumerate(langs):
    y = 3.45 + i * 0.62
    text(s, 1.4, y, 0.7, 0.4, lg, size=12, color=ORANGE, bold=True)
    text(s, 2.15, y, 4.9, 0.4, val, size=15, color=INK)
box(s, 7.75, 2.1, 4.6, 1.75, WHITE)
accent_bar(s, 7.75, 2.36, 1.23, ORANGE)
text(s, 8.1, 2.35, 4.0, 0.4, "The data translates", size=16.5, color=INK, bold=True)
text(s, 8.1, 2.85, 4.05, 0.8, "each language stored separately, side by side, on the same record",
     size=12, color=MUTED, spacing=1.25)
box(s, 7.75, 4.15, 4.6, 1.9, WHITE)
accent_bar(s, 7.75, 4.41, 1.38, TEAL)
text(s, 8.1, 4.4, 4.0, 0.4, "The interface translates", size=16.5, color=INK, bold=True)
text(s, 8.1, 4.9, 4.05, 0.9,
     "every person works in their own language — buttons, menus, labels",
     size=12, color=MUTED, spacing=1.25)
ctext(s, 6.5, "…and some values never translate: the inventory number, the coordinates, the year.",
      size=13.5, color=MUTED, italic=True)
footer(s, page=9)
notes(s, """
Multilingual is a word institutions use loosely, so let's be precise. In a
heritage system, two different things get translated, and confusing them
causes endless trouble.

First, the data. A title, a description, an abstract can exist in as many
languages as your institution works in — each version stored separately, on
the same record. The four lines you see here are one field, one record, four
languages. A researcher reads in English; a volunteer enters in Catalan; both
see the same object.

Second, the interface — the buttons, menus and field labels your staff see.
That translates too, independently. A translation workflow can run over your
data without ever touching the program.

And note what deliberately does NOT translate: the inventory number, a
coordinate, a year. Those are the same in every language, and Dédalo knows
the difference — one more consequence of fields that understand their kind.

For heritage data in particular — small languages, cross-border research,
diaspora audiences — this isn't a checkbox feature. It's how a collection
speaks to everyone it belongs to.
""")

# ── 10 · Time machine ──────────────────────────────────────────────────────
s = slide(PAPER, "Nothing is silently lost")
kicker(s, M, 0.62, "THE TIME MACHINE")
title(s, "Nothing is silently lost.", size=36)
box(s, 1.5, 2.3, 5.2, 2.75, C(0xEF, 0xEC, 0xE6), line=C(0xDD, 0xD7, 0xCC))
text(s, 1.85, 2.47, 4.6, 0.3, "v2 · 2019 · J. Ferrer", size=10.5, color=MUTED, bold=True)
box(s, 1.95, 2.85, 5.2, 2.75, WHITE, line=HAIR)
text(s, 2.3, 3.02, 4.6, 0.3, "v7 · 2023 · M. Puig", size=10.5, color=MUTED, bold=True)
box(s, 2.4, 3.4, 5.2, 2.75, WHITE, line=ORANGE)
text(s, 2.75, 3.57, 4.6, 0.3, "v12 · today · A. Serra", size=10.5, color=ORANGE,
     bold=True)
text(s, 2.75, 3.98, 4.5, 0.35, "Coins · inv. 1234", size=13, color=INK, bold=True)
bar(s, 2.75, 4.45, 4.3, 0.2)
bar(s, 2.75, 4.85, 3.2, 0.2)
bar(s, 2.75, 5.25, 3.8, 0.2)
chip(s, 2.75, 5.62, 3.1, 0.42, "restore any earlier state", fill=TINT[ORANGE],
     line_c=ORANGE, color=ORANGE, size=11, bold=True)
tm = [
    ("Every edit is kept", "what it was, who changed it, when — all of it", ORANGE),
    ("History you can read", "open any record and follow its whole life", TEAL),
    ("Recover anything", "a wrong edit, a bad bulk change, a deletion", VIOLET),
]
for i, (head, sub, col) in enumerate(tm):
    card(s, 8.15, 2.1 + i * 1.5, 4.25, 1.35, head, [sub], color=col, hsize=15.5,
         bsize=11.5)
ctext(s, 6.6, "This is what makes an attribution defensible and a correction traceable.",
      size=13.5, color=INK, italic=True)
footer(s, page=10)
notes(s, """
Every cataloguer has lived the moment: someone changed a description, and
nobody knows what it said before, or why.

In Dédalo, nothing is silently lost. Every change to a record is kept — what
it was, who made the change, and when. The tool is called the time machine,
and it does exactly what it says: you open a record and read its biography,
version by version, and you can restore any earlier state.

For a scientific catalogue this is not a comfort feature. Think what it
actually guarantees. An attribution is defensible — you can show the evidence
as it stood on the day you made the call. A correction is traceable — you can
prove you fixed it, and when. A disputed record is reconstructible — even
after a bad bulk change by a tired intern.

And there's a human consequence that matters as much: you can hand a
collection to a colleague — a successor, a partner institution — and they
receive not just the answers but the reasoning that produced them.

Twenty years of decisions, kept. That's what "for decades" meant on slide
three.
""")

# ── 11 · Media ─────────────────────────────────────────────────────────────
s = slide(PAPER, "Files are handled like heritage")
kicker(s, M, 0.62, "FILES, HANDLED LIKE HERITAGE")
title(s, "The master is never touched.", size=36)
mb = box(s, 1.05, 2.5, 2.7, 1.7, DEEP2)
text(s, 1.35, 2.75, 2.2, 0.3, "MASTER", size=11, color=ORANGE_L, bold=True,
     tracking=180)
text(s, 1.35, 3.12, 2.3, 0.4, "entrevista_0204.mov", size=11.5, color=PAPER)
text(s, 1.35, 3.55, 2.3, 0.4, "never altered — ever", size=10.5, color=MUTED_D,
     italic=True)
outs = [("Derivatives", 2.35), ("Thumbnails", 3.15), ("Transcript", 3.95)]
for lbl, y in outs:
    line(s, 3.75, 3.35, 4.6, y + 0.23, color=C(0xC9, 0xC2, 0xB4), w=1.2)
    chip(s, 4.6, y, 2.0, 0.46, lbl, fill=WHITE, line_c=HAIR, color=INK, size=12,
         bold=True)
import math
wx = 1.15
for i in range(26):
    hh = 0.12 + 0.30 * abs(math.sin(i * 0.55))
    col = ORANGE if i == 11 else C(0xD8, 0xD3, 0xC8)
    bar(s, wx, 5.75 - hh / 2, 0.06, hh, fill=col)
    wx += 0.21
text(s, 3.4, 6.05, 2.6, 0.3, "12:34 — a citable passage", size=10.5, color=ORANGE,
     bold=True)
mm = [
    ("Big files are ordinary", "a 32 GB video is a normal day — served directly, not pushed through the program", ORANGE),
    ("Passages are citable", "a line of testimony, with its timecode; a paragraph, with its page reference", TEAL),
]
for i, (head, sub, col) in enumerate(mm):
    card(s, 7.35, 2.35 + i * 1.75, 5.05, 1.55, head, [sub], color=col, hsize=16.5,
         bsize=12)
footer(s, page=11)
notes(s, """
Heritage collections are no longer just text. Photographs, recordings, film,
scanned documents — and Dédalo treats files the way a conservation lab treats
objects, not the way an email client treats attachments.

When you upload something, it becomes the master — the conserved original.
Dédalo never alters it. Ever. Everything the interface needs — smaller
versions, thumbnails, formats for the web — are generated as derivatives
beside it. The original stays exactly as you deposited it, thirty years from
now.

For oral history there's a second layer: recordings can be transcribed and
indexed, so that a passage — with its timecode — becomes a citable unit. A
researcher doesn't cite "interview 204, somewhere around minute twelve".
They cite 12:34, and the system takes them there. Long texts work the same
way with page references.

And scale: large files are served directly by the web server rather than
carried through the program, which is why a multi-gigabyte video is an
ordinary record, not an incident.

The principle is conservation: handle bytes like heritage.
""")

# ── 12 · Publishing ────────────────────────────────────────────────────────
s = slide(PAPER, "Publication is a decision")
kicker(s, M, 0.62, "PUBLISHING")
title(s, "Publication is a decision.", size=36)
recs = ["✓  Coin · 1234", "✓  Hoard · La Cova", "✓  Site · Tossal Rodó",
        "—   Donor file · XX", "—   Unpublished · 88"]
for i, r in enumerate(recs):
    pub = r.startswith("✓")
    chip(s, 1.05, 2.3 + i * 0.62, 2.5, 0.48, r, fill=(TINT[ORANGE] if pub else WHITE),
         line_c=(ORANGE if pub else HAIR), color=(ORANGE if pub else MUTED_D),
         size=10.5, bold=pub)
arrow_h(s, 3.75, 3.85, 0.75, color=ORANGE)
chip(s, 4.6, 3.45, 1.9, 0.8, "PUBLISH", fill=VIOLET, line_c=None, color=WHITE,
     size=14, bold=True)
text(s, 4.35, 4.32, 2.4, 0.3, "automatic · resumable", size=9.5, color=MUTED,
     align=PP_ALIGN.CENTER)
arrow_h(s, 6.6, 3.85, 0.75, color=ORANGE)
box(s, 7.45, 3.3, 2.6, 1.1, TINT[TEAL], line=TEAL)
text(s, 7.45, 3.5, 2.6, 0.3, "PUBLIC COPY", size=10.5, color=TEAL, bold=True,
     align=PP_ALIGN.CENTER, tracking=150)
text(s, 7.45, 3.85, 2.6, 0.3, "generated · read-only", size=11, color=MUTED,
     align=PP_ALIGN.CENTER)
arrow_h(s, 10.15, 3.85, 0.75, color=ORANGE)
text(s, 11.0, 3.55, 1.6, 0.7, ["your website", "open data"], size=11.5, color=INK,
     bold=True, align=PP_ALIGN.CENTER, spacing=1.2)
pubs = [
    ("Record by record", "no switch ever makes a whole collection public", ORANGE),
    ("Reversible", "un-mark a record; the next run removes it", TEAL),
    ("Field by field", "public record, private donor name", VIOLET),
]
for i, (head, sub, col) in enumerate(pubs):
    card(s, M + i * 3.99, 5.35, 3.7, 1.35, head, [sub], color=col, tint=True,
         hsize=15.5, bsize=11.5)
footer(s, page=12)
notes(s, """
We've built good data, in private, in the workshop. Now: the world. And the
governing principle is that publication is a decision — never an accident.

There is no "make the archive public" switch in Dédalo. You mark records for
publication, one by one. A background job then copies the marked records out
into the public database, and writes whatever file formats your partners and
aggregators want — CSV, XML, RDF. The job runs on its own, and if it's
interrupted, it resumes where it stopped.

Three properties to notice. Per record: what you didn't mark simply doesn't
exist out there. Reversible: un-mark a record and the next run removes it —
taking something down is as easy as putting it up. And per field: a record can
be public while parts of it aren't — the donor's name, a sensitive location,
personal data of living people, testimony given under condition.

Why does this design matter? Because it's what lets you catalogue frankly —
with internal notes and unresolved attributions — and still publish
confidently. A collection can be catalogued for years and published in one
afternoon.
""")

# ── 13 · Three doors ───────────────────────────────────────────────────────
s = slide(PAPER, "Three doors, one lock")
kicker(s, M, 0.62, "THREE DOORS INTO THE SAME COLLECTION")
title(s, "Three doors. One lock.", size=36)
doors = [
    ("STAFF", ORANGE, "Your team, signed in",
     "everything their permissions allow — set field by field, section by section",
     "never lets anyone in without an account"),
    ("PUBLIC", TEAL, "The world reads the copy",
     "records, images, search — passages with their page or timecode",
     "never writes anything, anywhere"),
    ("AI", VIOLET, "An assistant, on a leash",
     "sees exactly what its user sees; when it wants to change something, it writes a plan for a human",
     "never saves an edit nobody confirmed"),
]
for i, (name, col, who, may, never) in enumerate(doors):
    x = M + i * 3.99
    box(s, x, 2.05, 3.7, 4.0, WHITE, line=C(0xD4, 0xCE, 0xC2))
    disc(s, x + 1.85, 2.62, 0.66, col)
    text(s, x, 3.05, 3.7, 0.35, name, size=12.5, color=col, bold=True,
         align=PP_ALIGN.CENTER, tracking=200)
    text(s, x + 0.3, 3.45, 3.1, 0.4, who, size=14.5, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, x + 0.3, 4.05, 3.1, 1.0, may, size=11.5, color=MUTED,
         align=PP_ALIGN.CENTER, spacing=1.25)
    rule(s, x + 1.35, 5.15, 1.0, color=HAIR)
    text(s, x + 0.3, 5.3, 3.1, 0.65, never, size=11, color=ROSE, bold=True,
         align=PP_ALIGN.CENTER, spacing=1.2)
ctext(s, 6.45, "Whoever knocks — a person, a website, an assistant — the same checks run, in the same order.",
      size=13.5, color=INK, italic=True)
footer(s, page=13)
notes(s, """
You'll hear the word "API" in every technical conversation about Dédalo. An
API is simply a door other programs may knock on. Dédalo has three, and this
slide is really about trust.

The staff door: your team, each with a personal account, doing everything
their permissions allow. And permissions are granular — per section, per
field. A volunteer can catalogue physical descriptions while never seeing the
donor's name. That's not a policy you hope people follow; it's the shape of
the system.

The public door: your website and the world. It reads the published copy and
it can only hand things out — the door has no handle on the inside.

The AI door: an assistant can search and read — seeing exactly what its user
sees, never more — and when it wants to change something, it writes a plan,
op by op, that a human confirms. It cannot save on its own.

Now the point directors care about: three doors, one lock. Whoever knocks —
person, website, or assistant — passes the same checks, in the same order,
against the same permissions. No back doors, no shortcuts. Opening the AI
door does not open anything else.
""")

# ── 14 · AI / semantic search ──────────────────────────────────────────────
s = slide(PAPER, "Ready for AI, on your terms")
kicker(s, M, 0.62, "READY FOR AI — ON YOUR TERMS")
title(s, "Find it by meaning.", size=36)
box(s, 1.05, 2.15, 6.1, 1.0, TINT[ROSE], line=ROSE)
text(s, 1.35, 2.32, 5.6, 0.7,
     "“Testimonies about people who had to leave their homes when the dam was built.”",
     size=14, color=INK, italic=True, spacing=1.2)
matches = [
    ("“…when the water came and we had to leave…”", "Testimony 204 · 07:12"),
    ("“…el pantano nos dejó sin casa…”", "Testimony 118 · 22:41"),
    ("“…they flooded our village, that's how we lost it…”", "Testimony 231 · 03:55"),
]
for i, (q, cite) in enumerate(matches):
    y = 3.45 + i * 0.82
    box(s, 1.05, y, 6.1, 0.68, WHITE, line=HAIR)
    text(s, 1.3, y + 0.1, 4.3, 0.4, q, size=12.5, color=INK)
    text(s, 5.35, y + 0.18, 1.7, 0.35, cite, size=9.5, color=TEAL, bold=True,
         align=PP_ALIGN.RIGHT)
text(s, 1.05, 6.0, 6.1, 0.35,
     "a word search finds none of these — a dam is not a pantano is not “the water”",
     size=11, color=MUTED, italic=True)
ai = [
    ("Asks, never takes", "your data is not handed to a model to keep", ORANGE),
    ("Sees only what you see", "an assistant works at its user's permission level", TEAL),
    ("Proposes, you decide", "it cannot save an edit; a person confirms", VIOLET),
]
for i, (head, sub, col) in enumerate(ai):
    card(s, 7.65, 2.15 + i * 1.42, 4.75, 1.28, head, [sub], color=col, hsize=15.5,
         bsize=11.5)
text(s, 7.65, 6.45, 4.75, 0.5,
     "Restricted collections can be locked inside the building — even from AI.",
     size=11.5, color=ROSE, bold=True, spacing=1.2)
footer(s, page=14)
notes(s, """
Two AI pieces are built into Dédalo. Both are optional, both are off until
your administrator says otherwise — and both are shaped for heritage work.

First: semantic search — finding things by what they mean, not by exact
words. A researcher asks about people displaced by a dam. Word search finds
only records where someone wrote "dam". Dédalo also finds "when the water
came and we had to leave", "el pantano", "they flooded our village" — because
those mean almost the same thing. And every result cites the record it came
from, so a researcher can open it and judge for themselves.

Why does heritage need this more than most data? Because heritage data is
multilingual, historical, dialectal, paraphrastic — testimony spoken in 1960s
vocabulary simply doesn't contain the researcher's keywords.

Note what it does not replace: precise field search stays — "every coin
minted before 100 BC" is still exact. Meaning is an additional door.

Second: AI assistants plug in through the door we saw last slide. It asks
Dédalo — never the other way round. It sees only what its user sees. And it
proposes; a person decides. Sensitive collections can be excluded entirely,
or run against a local model so nothing leaves the building.
""")

# ── 15 · Trust ─────────────────────────────────────────────────────────────
s = slide(PAPER, "Twenty years in the field, open by design")
kicker(s, M, 0.62, "WHY TRUST IT")
title(s, "Twenty years in the field. Open by design.", size=34)
text(s, M, 2.05, 5.6, 0.35, "IN PRODUCTION AT", size=10.5, color=MUTED_D, bold=True,
     tracking=180)
inst = ["Moneda Ibérica", "MUPREVA · València", "Memorial Democràtic",
        "Freie Universität Berlin", "Museu de la Paraula", "Lur Azpian · Navarra",
        "Mujer y Memoria", "MUVAET · Etnologia"]
for i, t in enumerate(inst):
    x = M + (i % 2) * 2.95
    y = 2.5 + (i // 2) * 0.72
    chip(s, x, y, 2.75, 0.52, t, fill=WHITE, line_c=HAIR, color=INK, size=11)
text(s, M, 5.55, 5.9, 0.4, "coins · archaeology · oral memory · ethnology · exhumations",
     size=11.5, color=MUTED, italic=True)
trust = [
    ("Free and open source", "no licence, no vendor, no lock-in — the code is yours to read"),
    ("Your data, in the open", "PostgreSQL; exports to RDF, Dublin Core, JSON-LD, CSV, XML"),
    ("Quality you can measure", "“Raspa” — a 0-to-10 data-quality score for your catalogue"),
]
for i, (head, sub) in enumerate(trust):
    y = 2.25 + i * 1.32
    ring(s, 7.85, y + 0.42, 0.5, GREEN, "✓", color=GREEN, size=15)
    text(s, 8.35, y + 0.12, 4.1, 0.4, head, size=16.5, color=INK, bold=True)
    text(s, 8.35, y + 0.55, 4.05, 0.6, sub, size=11.5, color=MUTED, spacing=1.2)
ctext(s, 6.45, "The engine was rebuilt in 2026. The data was never touched.",
      size=14.5, color=INK, bold=True)
footer(s, page=15)
notes(s, """
Why should an institution trust this with its memory? Four reasons.

It is in production, today, in the field — not a prototype looking for users.
Numismatic catalogues in Spain, ethnological collections in València, oral
memory archives from the Spanish Civil War, archaeological and exhumation
records, research at the Freie Universität Berlin. Coins and testimony,
twenty years of accumulated practice.

It is free and open source. No licence fees, no vendor who can disappear or
reprice you, no format held hostage. The code is there to be read — which,
for an institution whose duty is preservation, is itself a preservation
strategy.

Your data is yours, in PostgreSQL, and it speaks the open languages of the
field: RDF, Dublin Core, JSON-LD, CSV, XML. Aggregators and partners can
harvest it through the public door on standard terms.

And quality is measurable: Dédalo carries a score called Raspa — zero to ten,
from "structured" up through "traceable, translatable, openly processable".
Useful for arguing with evidence in a funding application.

Last line: in 2026 the engine behind it was completely rebuilt — and not one
record, not one byte of anyone's data was touched. That's the durability test
it chose to pass.
""")

# ── 16 · Closing ───────────────────────────────────────────────────────────
s = slide(DEEP, "Memory, kept alive")
for (x1, y1, x2, y2) in [(10.6, 1.2, 11.9, 2.0), (11.9, 2.0, 11.1, 3.1),
                          (11.1, 3.1, 12.4, 3.9)]:
    line(s, x1, y1, x2, y2, color=LINK_D, w=1.2)
disc(s, 10.6, 1.2, 0.3, TEAL)
disc(s, 11.9, 2.0, 0.45, ORANGE)
disc(s, 11.1, 3.1, 0.24, ORANGE_L)
disc(s, 12.4, 3.9, 0.3, VIOLET)
text(s, M, 1.35, 11.0, 1.2, "Memory, kept alive.", size=54, color=PAPER, bold=True)
rule(s, M, 2.75, 2.3, color=ORANGE_L, thickness=0.03)
text(s, M, 3.15, 10.5, 0.4, "IF YOU REMEMBER THREE THINGS:", size=12, color=ORANGE_L,
     bold=True, tracking=200)
three = [
    "Your catalogue is private, and it is the original.",
    "Its shape is yours: fields, records and vocabularies are data you edit.",
    "Nothing happens silently — every edit is kept, every record published by decision.",
]
for i, t in enumerate(three):
    y = 3.75 + i * 0.62
    disc(s, M + 0.18, y + 0.16, 0.34, ORANGE, str(i + 1), size=12)
    text(s, M + 0.6, y, 10.5, 0.4, t, size=16, color=PAPER)
text(s, M, 6.0, 8.0, 0.4, "dedalo.dev   ·   demo.dedalo.dev", size=15, color=ORANGE_L,
     bold=True)
text(s, M, 6.5, 8.0, 0.4, "Questions — and a live look at a real catalogue.", size=13,
     color=MUTED_D)
footer(s, "Dédalo · dedalo.dev", color=C(0x5C, 0x6A, 0x7C))
notes(s, """
Let's close the way the project's own manual closes, with three things to
remember.

One: your catalogue is private, and it is the original. What the public sees
is a copy, on a different machine, that can only be read. Nothing is public
by accident.

Two: the shape of the catalogue is yours. Fields, record types and
vocabularies are data you edit inside Dédalo — not code somebody has to
rewrite for you. Your model can follow your discipline, and evolve with it.

Three: nothing is published or changed silently. Publication is a decision,
record by record, and every edit ever made is kept, with its author and its
date.

That's the whole story, really: produce good data, keep it for decades,
publish by decision — in a workshop whose shape belongs to the people who
know the collection.

Thank you. The demo is live and open — let's go look at a real catalogue and
see the time machine and the thesaurus in action, and I'll take your
questions while we do it.
""")

# ═════════════════════════════════════════════════════════════════════════
#  OUTPUT
# ═════════════════════════════════════════════════════════════════════════
PPTX = os.path.join(OUT_DIR, "dedalo_for_curators.pptx")
prs.save(PPTX)

with open(os.path.join(OUT_DIR, "speaker_script.md"), "w") as f:
    f.write("# Dédalo for curators — speaker script\n\n")
    f.write("Companion to `dedalo_for_curators.pptx` (the same text is embedded as "
            "notes). 16 slides, ~2,700 words ≈ 21 min of talk — leaves ~10 min for "
            "questions and the live look. Slides 5, 7, 13 and 14 carry the argument; "
            "if pressed for time, skim 9 and 11.\n\n")
    for i, (t, n) in enumerate(SCRIPT, 1):
        f.write(f"## Slide {i} — {t}\n\n{n}\n\n")

words = sum(len(n.split()) for _, n in SCRIPT)
print(f"wrote {PPTX} ({len(prs.slides._sldIdLst)} slides)  "
      f"script: {words} words ≈ {words/130:.0f} min at 130 wpm")
print(f"wrote {os.path.join(OUT_DIR, 'speaker_script.md')}")
